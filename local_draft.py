# Import necessary libraries
import os
import faiss
import fitz  # PyMuPDF
from typing import List, Dict, Optional, Tuple
from datetime import datetime
import requests
import numpy as np
from PIL import Image
import io
import base64
import cv2
import pytesseract
from difflib import SequenceMatcher
from config import OPENAI_API_KEY
import logging
import time
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
import openai
import pdfplumber
from abc import ABC, abstractmethod
from docx import Document as DocxDocument
from pptx import Presentation
import streamlit as st
import openpyxl
import pandas as pd
import re
import tiktoken
from duckduckgo_search import DDGS
import speech_recognition as sr
from pydub import AudioSegment
import whisper
import tempfile
import subprocess
from pdf2image import convert_from_path
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import black, blue
from reportlab.lib.enums import TA_LEFT, TA_CENTER

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Reduce verbose httpx logging
logging.getLogger("httpx").setLevel(logging.WARNING)

### =================== Text Processing =================== ###
class TextProcessor:
    def __init__(self, chunk_size: int = 2000, overlap: int = 100):
        self.chunk_size = chunk_size
        self.overlap = overlap
        self.extracted_images = {}  # Store images for display
        self.image_descriptions = {}  # Store semantic descriptions for search
        self.images_analyzed = set()  # Track which images have been analyzed
        self.processing_lock = threading.Lock()  # Thread safety for image processing
        self.error_count = 0  # Track errors for better error handling

    def extract_text_from_pdf(self, pdf_path: str, enable_image_processing: bool = False) -> str:
        """Extract text and images from PDF, detect tables/graphs, and run OCR on each region (header + body)."""
        try:
            doc = fitz.open(pdf_path)
            text_content = []
            
            # Disable image processing for now
            enable_image_processing = False
            
            if enable_image_processing:
                os.makedirs("images", exist_ok=True)
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                text_content.append(f"Page {page_num + 1}: {text}")

                # Only process images if enabled (much faster without this)
                if enable_image_processing:
                    try:
                        # Render page as image with lower DPI for speed
                        pix = page.get_pixmap(dpi=150)  # Reduced from 200 to 150
                        img_data = pix.tobytes("png")
                        img_array = np.frombuffer(img_data, np.uint8)
                        
                        # Check image size to prevent processing extremely large images
                        if len(img_array) > 25 * 1024 * 1024:  # Reduced to 25MB limit
                            continue
                        
                        # Add error handling for large PNG chunks
                        img_cv = cv2.imdecode(img_array, cv2.IMREAD_COLOR)
                        if img_cv is not None:
                            detected_regions = self.detect_tables_and_graphs(img_cv)
                            for idx, (x, y, w, h) in enumerate(detected_regions):
                                crop = img_cv[y:y+h, x:x+w]
                                img_filename = f"page_{page_num + 1}_detected_{idx + 1}.png"
                                img_path = os.path.join("images", img_filename)
                                cv2.imwrite(img_path, crop)
                                img_key = f"page_{page_num + 1}_detected_{idx + 1}"
                                # OCR on the cropped region (body)
                                ocr_text = self.ocr_image(crop)
                                # OCR on the header (top 20%)
                                header_h = max(1, int(h * 0.2))
                                header_crop = crop[0:header_h, :]
                                ocr_header = self.ocr_image(header_crop)
                                self.extracted_images[img_key] = {
                                    'path': img_path,
                                    'page': page_num + 1,
                                    'image_num': idx + 1,
                                    'description': 'Detected table/graph region (OpenCV)',
                                    'ocr_text': ocr_text,
                                    'ocr_header': ocr_header,
                                }
                    except Exception as e:
                        logger.warning(f"Error processing images on page {page_num + 1}: {str(e)}")
                        continue
            
            doc.close()
            final_content = "\n".join(text_content)
            return final_content
        except Exception as e:
            self.error_count += 1
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return f"Error extracting text from PDF: {str(e)}"

    def extract_text_from_docx(self, path):
        try:
            doc = DocxDocument(path)
            return '\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
        except Exception as e:
            return f"Error extracting text from Word: {str(e)}"

    def extract_text_from_pptx(self, path):
        try:
            prs = Presentation(path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return '\n'.join([t for t in text_runs if t.strip()])
        except Exception as e:
            return f"Error extracting text from PowerPoint: {str(e)}"

    def extract_text_from_xlsx(self, path):
        """Extract comprehensive, queryable content from Excel files with proper company-by-company analysis."""
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            output = []
            
            for ws in wb.worksheets:
                output.append(f"\n=== Sheet: {ws.title} ===\n")
                
                # Get all rows as lists
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    output.append("(Sheet is empty)\n")
                    continue
                
                # Use first row as header
                headers = [str(cell) if cell is not None else f"Column_{i+1}" for i, cell in enumerate(rows[0])]
                output.append(f"Columns: {', '.join(headers)}\n")
                output.append(f"Total Rows: {len(rows)-1}\n")
                
                # Process data rows with SAFE LIMITS
                data_rows = rows[1:]
                MAX_ROWS = 1000  # Safe limit to prevent crashes
                
                if len(data_rows) > MAX_ROWS:
                    output.append(f"\nWARNING: Processing first {MAX_ROWS} rows out of {len(data_rows)} to prevent memory issues.\n")
                    data_rows = data_rows[:MAX_ROWS]
                
                # === COMPANY-BY-COMPANY ANALYSIS FORMAT ===
                output.append("\n=== INDIVIDUAL COMPANY ANALYSIS ===\n")
                
                # Process each company/row individually
                for row_idx, row in enumerate(data_rows, 1):
                    output.append(f"\n--- COMPANY {row_idx} ---\n")
                    
                    # Create a structured company profile
                    company_data = {}
                    for col_idx, cell in enumerate(row):
                        if cell is not None:
                            # Format numbers nicely
                            if isinstance(cell, (int, float)):
                                if isinstance(cell, int):
                                    formatted_cell = f"{cell:,}"
                                else:
                                    formatted_cell = f"{cell:.2f}" if cell != int(cell) else f"{cell:.0f}"
                            else:
                                formatted_cell = str(cell)
                            company_data[headers[col_idx]] = formatted_cell
                        else:
                            company_data[headers[col_idx]] = "(empty)"
                    
                    # Output company data in a clear, structured format
                    for header, value in company_data.items():
                        output.append(f"{header}: {value}\n")
                    
                    # Add a separator for readability
                    output.append("-" * 50 + "\n")
                
                # Add summary statistics
                output.append("\n=== SUMMARY STATISTICS ===\n")
                
                # Count by key metrics
                for col_idx, header in enumerate(headers):
                    if any('name' in header.lower() or 'company' in header.lower() for word in header.split()):
                        non_empty = sum(1 for row in data_rows if row[col_idx] is not None and str(row[col_idx]).strip())
                        output.append(f"Companies with {header}: {non_empty}\n")
                    
                    if any('nav' in header.lower() or 'value' in header.lower() or 'revenue' in header.lower() for word in header.split()):
                        numeric_values = []
                        for row in data_rows:
                            if row[col_idx] is not None and isinstance(row[col_idx], (int, float)):
                                numeric_values.append(row[col_idx])
                        
                        if numeric_values:
                            output.append(f"{header} - Total: {sum(numeric_values):,.2f}, Avg: {sum(numeric_values)/len(numeric_values):,.2f}, Max: {max(numeric_values):,.2f}\n")
                
                # Top companies by key metrics
                output.append("\n=== TOP COMPANIES ===\n")
                
                # Find NAV or value columns
                value_columns = []
                for col_idx, header in enumerate(headers):
                    if any('nav' in header.lower() or 'value' in header.lower() or 'revenue' in header.lower() for word in header.split()):
                        value_columns.append((col_idx, header))
                
                if value_columns:
                    for col_idx, header in value_columns:
                        # Get top 10 companies by this metric
                        companies_with_values = []
                        for row_idx, row in enumerate(data_rows, 1):
                            if row[col_idx] is not None and isinstance(row[col_idx], (int, float)):
                                # Try to find company name
                                company_name = "Unknown"
                                for name_col_idx, name_header in enumerate(headers):
                                    if any('name' in name_header.lower() or 'company' in name_header.lower() for word in name_header.split()):
                                        if row[name_col_idx] is not None:
                                            company_name = str(row[name_col_idx])
                                            break
                                
                                companies_with_values.append((company_name, row[col_idx], row_idx))
                        
                        # Sort by value and show top 10
                        companies_with_values.sort(key=lambda x: x[1], reverse=True)
                        output.append(f"\nTop 10 companies by {header}:\n")
                        for i, (name, value, row_num) in enumerate(companies_with_values[:10], 1):
                            formatted_value = f"{value:,.2f}" if isinstance(value, (int, float)) else str(value)
                            output.append(f"{i}. {name} (Row {row_num}): {formatted_value}\n")
            
            return "\n".join(output)
            
        except Exception as e:
            # Fallback to pandas with same structured format
            try:
                df = pd.read_excel(path, engine='openpyxl')
                output = [f"Sheet: {getattr(df, 'sheet_name', 'N/A')} (via pandas)\n"]
                output.append(f"Shape: {df.shape}\n")
                output.append(f"Columns: {list(df.columns)}\n")
                
                # Limit rows for safety
                MAX_ROWS = 1000
                if len(df) > MAX_ROWS:
                    output.append(f"\nWARNING: Processing first {MAX_ROWS} rows out of {len(df)} to prevent memory issues.\n")
                    df = df.head(MAX_ROWS)
                
                output.append("\n=== INDIVIDUAL COMPANY ANALYSIS ===\n")
                
                # Process each company individually
                for idx, row in df.iterrows():
                    output.append(f"\n--- COMPANY {idx+1} ---\n")
                    for col, val in row.items():
                        if pd.notna(val):
                            if isinstance(val, (int, float)):
                                formatted_val = f"{val:,.2f}" if val != int(val) else f"{val:,}"
                            else:
                                formatted_val = str(val)
                            output.append(f"{col}: {formatted_val}\n")
                        else:
                            output.append(f"{col}: (empty)\n")
                    output.append("-" * 50 + "\n")
                
                return "\n".join(output)
                
            except Exception as e2:
                return f"Error extracting text from Excel: {str(e)}; {str(e2)}"

    def extract_text_from_xls(self, path):
        try:
            df = pd.read_excel(path)
            return df.to_string(index=False)
        except Exception as e:
            return f"Error extracting text from Excel: {str(e)}"

    def detect_tables_and_graphs(self, img_cv):
        """Detect likely tables/graphs in a page image using OpenCV (returns list of bounding boxes)."""
        try:
            gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
            blur = cv2.GaussianBlur(gray, (5, 5), 0)
            edges = cv2.Canny(blur, 50, 150, apertureSize=3)
            # Find contours
            contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            regions = []
            for cnt in contours:
                x, y, w, h = cv2.boundingRect(cnt)
                # Heuristic: Only keep large regions (likely tables/graphs)
                if w > 100 and h > 60 and w*h > 10000:
                    regions.append((x, y, w, h))
            # Optionally, merge overlapping/close rectangles here
            return regions
        except Exception as e:
            logger.warning(f"Error detecting tables/graphs: {str(e)}")
            return []

    def ocr_image(self, img_cv):
        """Run OCR on a cropped image region using Tesseract if available, else fallback to empty string."""
        try:
            # Convert to RGB for pytesseract
            img_rgb = cv2.cvtColor(img_cv, cv2.COLOR_BGR2RGB)
            text = pytesseract.image_to_string(img_rgb)
            return text.strip()
        except Exception as e:
            logger.warning(f"OCR error: {str(e)}")
            return ""

    def ensure_images_analyzed(self):
        """Analyze all images that have not yet been analyzed with Vision API."""
        with self.processing_lock:
            for img_key, img_info in list(self.extracted_images.items()):
                if img_key in self.images_analyzed:
                    continue
                try:
                    img_path = img_info['path']
                    if os.path.exists(img_path):
                        img_pil = Image.open(img_path)
                        img_analysis = self.analyze_image_detailed(img_pil)
                        if img_analysis:
                            self.image_descriptions[img_key] = img_analysis
                        else:
                            # If not a figure/graph, remove from extracted images
                            del self.extracted_images[img_key]
                            if os.path.exists(img_path):
                                os.remove(img_path)
                    self.images_analyzed.add(img_key)
                except Exception as e:
                    logger.warning(f"Error analyzing image {img_key}: {str(e)}")
                    continue

    def analyze_image_detailed(self, img_pil):
        """Enhanced image analysis using Vision API for semantic search - focused on figures and graphs."""
        try:
            # Resize image to reduce API costs and speed
            max_size = 512
            img_pil.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
            
            # Convert to base64
            buffer = io.BytesIO()
            img_pil.save(buffer, format='PNG')
            img_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            # Enhanced Vision API call focused on figures and graphs
            headers = {
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Analyze this image and determine if it contains figures, graphs, charts, tables, or data visualizations. If it does, provide a detailed description including: 1) Type of visual (chart, graph, table, diagram, etc.) 2) Main subject/topic 3) Key data points or information shown 4) Any text or labels visible 5) Business context if applicable (revenue, growth, metrics, etc.). If this is NOT a figure/graph/chart (e.g., it's just text, a photo, or decorative element), respond with 'NOT_A_FIGURE'. Be specific and searchable."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}"
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 16384
            }
            
            response = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers=headers,
                json=payload,
                timeout=100
            )
            
            if response.status_code == 200:
                result = response.json()
                content = result["choices"][0]["message"]["content"]
                
                # Only return description if it's actually a figure/graph
                if "NOT_A_FIGURE" not in content.upper():
                    return content
                else:
                    return None
            else:
                return None
                
        except Exception as e:
            logger.warning(f"Vision API error: {str(e)}")
            return None

    def analyze_image_simple(self, img_pil):
        """Simple image analysis using Vision API - minimal and fast."""
        try:
            # Resize image to reduce API costs and speed
            max_size = 512
            img_pil.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
            
            # Convert to base64
            buffer = io.BytesIO()
            img_pil.save(buffer, format='PNG')
            img_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            # Simple Vision API call
            headers = {
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Describe this image briefly in one sentence. Focus on charts, tables, graphs, or key visual elements."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{img_base64}"
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 16384
            }
            
            response = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers=headers,
                json=payload,
                timeout=100
            )
            
            if response.status_code == 200:
                result = response.json()
                return result["choices"][0]["message"]["content"]
            else:
                return None
                
        except Exception as e:
            logger.warning(f"Simple Vision API error: {str(e)}")
            return None

    def fuzzy_score(self, a, b):
        """Return a fuzzy match ratio between two strings (0-100)."""
        return int(SequenceMatcher(None, a, b).ratio() * 100)

    def search_images_semantically(self, query: str, top_k: int = 3, min_score: int = 180):
        """Find the most relevant detected region by fuzzy matching query to header/body OCR, prioritizing header. Only return if score is high enough."""
        try:
            query_lc = query.lower()
            
            # Check if this is a general image request (not specific)
            general_image_requests = ['image', 'picture', 'show me', 'display', 'give me', 'single image', 'any image']
            is_general_request = any(phrase in query_lc for phrase in general_image_requests)
            
            # If it's a general request, return any available images
            if is_general_request:
                all_images = list(self.extracted_images.values())
                if all_images:
                    return all_images[:top_k]  # Return first few images
                else:
                    return []
            
            # For specific requests, use semantic matching
            best_score = -1
            best_img_info = None
            for img_info in self.extracted_images.values():
                ocr_header = img_info.get('ocr_header', '').lower()
                ocr_text = img_info.get('ocr_text', '').lower()
                header_score = self.fuzzy_score(query_lc, ocr_header)
                body_score = self.fuzzy_score(query_lc, ocr_text)
                combined_score = header_score * 2 + body_score  # Prioritize header
                if combined_score > best_score:
                    best_score = combined_score
                    best_img_info = img_info
            if best_img_info and best_score >= min_score:
                return [best_img_info]
            return []
        except Exception as e:
            logger.error(f"Error in semantic image search: {str(e)}")
            return []

    def is_blank_image(self, img_pil):
        """Check if image is mostly blank/white space."""
        try:
            # Convert to grayscale for easier processing
            gray_img = img_pil.convert('L')
            
            # Get pixel data
            pixels = list(gray_img.getdata())
            
            # Count non-white pixels (assuming white is close to 255)
            white_threshold = 250  # Very close to white
            non_white_pixels = sum(1 for pixel in pixels if pixel < white_threshold)
            
            # Calculate percentage of non-white pixels
            total_pixels = len(pixels)
            non_white_percentage = (non_white_pixels / total_pixels) * 100
            
            # Consider image blank if less than 5% non-white pixels
            return non_white_percentage < 5
            
        except Exception as e:
            return False  # If we can't check, assume it's not blank

    def get_image_path(self, page_num, image_num=None):
        """Get the path to a specific image for display."""
        if image_num:
            key = f"page_{page_num}_image_{image_num}"
        else:
            # Find first image on the page
            key = None
            for k in self.extracted_images.keys():
                if k.startswith(f"page_{page_num}_"):
                    key = k
                    break
        
        if key and key in self.extracted_images:
            return self.extracted_images[key]['path']
        return None

    def get_all_images(self):
        """Get all extracted images info."""
        return self.extracted_images

    def chunk_text(self, text: str) -> List[str]:
        """Split text into fixed-size chunks (no overlap, no sentence boundary logic)."""
        if not text.strip():
            return []
        chunk_size = self.chunk_size
        return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

    def prepare_documents(self, pdf_paths: List[str]) -> List[Dict]:
        """Prepare documents for vector storage."""
        documents = []
        
        for pdf_path in pdf_paths:
            try:
                # Extract text from PDF, but also get per-page text
                doc = fitz.open(pdf_path)
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    text = page.get_text()
                    if text.strip():
                        # Chunk this page's text
                        chunks = self.chunk_text(text)
                        for i, chunk in enumerate(chunks):
                            if chunk.strip():
                                documents.append({
                                    'text': chunk.strip(),
                                    'metadata': {
                                        'source': pdf_path,
                                        'chunk_id': i,
                                        'total_chunks': len(chunks),
                                        'date': datetime.now().strftime('%Y-%m-%d'),
                                        'page': page_num + 1
                                    }
                                })
                doc.close()
            except Exception as e:
                logger.error(f"Error preparing document {pdf_path}: {str(e)}")
                continue
        return documents

### =================== Document Loading =================== ###
class AudioProcessor:
    def __init__(self):
        self.ffmpeg_available = self._check_ffmpeg()
        self.whisper_model = None
        self.recognizer = sr.Recognizer()
        self._model_cache = {}  # Cache for different model sizes
        
        # Pre-load the tiny model for speed
        self.load_whisper_model("tiny")
    
    def _check_ffmpeg(self):
        """Check if FFmpeg is available on the system."""
        try:
            result = subprocess.run(['ffmpeg', '-version'], capture_output=True, text=True, timeout=10)
            if result.returncode == 0:
                logger.info("FFmpeg is available")
                return True
            else:
                logger.warning("FFmpeg check failed")
                return False
        except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
            logger.warning("FFmpeg not found or not accessible")
            return False
        except Exception as e:
            logger.error(f"Error checking FFmpeg: {str(e)}")
            return False
    
    def load_whisper_model(self, model_size="tiny"):
        """Load Whisper model with caching - use tiny for speed."""
        if model_size in self._model_cache:
            self.whisper_model = self._model_cache[model_size]
            return
        
        try:
            logger.info(f"Loading Whisper model: {model_size}")
            # Use tiny model for maximum speed (still very accurate)
            self.whisper_model = whisper.load_model(model_size)
            self._model_cache[model_size] = self.whisper_model
            logger.info("Whisper model loaded successfully")
        except Exception as e:
            logger.error(f"Error loading Whisper model: {str(e)}")
    
    def transcribe_with_whisper(self, audio_path: str) -> str:
        """Transcribe audio using Whisper with ULTRA speed optimizations."""
        try:
            if self.whisper_model is None:
                self.load_whisper_model("tiny")  # Use tiny model for speed
            
            logger.info(f"Transcribing audio with Whisper: {audio_path}")
            
            # Check if file exists and has content
            if not os.path.exists(audio_path):
                raise FileNotFoundError(f"Audio file not found: {audio_path}")
            
            file_size = os.path.getsize(audio_path)
            if file_size == 0:
                raise ValueError(f"Audio file is empty: {audio_path}")
            
            # Get audio duration
            import librosa
            duration = librosa.get_duration(path=audio_path)
            logger.info(f"Audio duration: {duration:.2f} seconds ({duration/60:.2f} minutes)")
            
            print(f"🎵 Transcribing {duration/60:.1f} minute audio file...")
            
            # ULTRA SPEED OPTIMIZED settings
            result = self.whisper_model.transcribe(
                audio_path,
                language="en",
                task="transcribe",
                verbose=False,  # Disable verbose for speed
                fp16=False,  # Disable FP16 for CPU compatibility and stability
                condition_on_previous_text=False,  # Disable for speed
                temperature=0.0,  # Deterministic for speed
                compression_ratio_threshold=2.0,  # More lenient for speed
                logprob_threshold=-3.0,  # More lenient for speed
                no_speech_threshold=0.9,  # More lenient for speed
                word_timestamps=False,  # Disable for speed
                prepend_punctuations=False,  # Disable for speed
                append_punctuations=False,  # Disable for speed
                initial_prompt=None,  # Disable for speed
                suppress_tokens=[-1],  # Suppress end token for speed
                without_timestamps=True  # Disable timestamps for speed
            )
            
            logger.info("Whisper transcription completed")
            
            transcription = result.get('text', '').strip()
            
            if not transcription:
                logger.warning("Whisper returned empty transcription")
                return "No speech detected in audio file."
            
            logger.info(f"Transcription completed: {len(transcription)} characters")
            print(f"🎉 Transcription completed! Total: {len(transcription)} characters")
            return transcription
            
        except Exception as e:
            # FIX: Capture the actual error details
            import traceback
            error_details = traceback.format_exc()
            logger.error(f"Error transcribing with Whisper: {str(e)}")
            logger.error(f"Full traceback: {error_details}")
            print(f"❌ Transcription error: {str(e)}")
            
            # If it's a tensor error, try with even more conservative settings
            if "tensor" in str(e).lower() and "reshape" in str(e).lower():
                print("🔄 Trying with ultra-conservative settings...")
                try:
                    result = self.whisper_model.transcribe(
                        audio_path,
                        language="en",
                        task="transcribe",
                        verbose=False,
                        fp16=False,  # Disable FP16 for stability
                        condition_on_previous_text=False,
                        temperature=0.0,
                        compression_ratio_threshold=2.0,
                        logprob_threshold=-3.0,
                        no_speech_threshold=0.9
                    )
                    transcription = result.get('text', '').strip()
                    if transcription:
                        return transcription
                except Exception as e2:
                    logger.error(f"Second attempt also failed: {str(e2)}")
            
            return f"Error transcribing audio: {str(e)}"
    
    def convert_audio_format(self, audio_path: str, target_format: str = "wav") -> str:
        """Convert audio to WAV format using FFmpeg with ULTRA optimized settings."""
        try:
            if not self.ffmpeg_available:
                logger.error("FFmpeg not available for audio conversion")
                return audio_path
            
            # Create output path
            base_name = os.path.splitext(audio_path)[0]
            output_path = f"{base_name}.{target_format}"
            
            logger.info(f"Converting {audio_path} to {output_path}")
            
            # ULTRA OPTIMIZED FFmpeg settings for maximum speed
            cmd = [
                'ffmpeg', '-i', audio_path,
                '-acodec', 'pcm_s16le',  # 16-bit PCM
                '-ar', '8000',           # 8kHz sample rate (faster processing)
                '-ac', '1',              # Mono
                '-y',                    # Overwrite output
                '-threads', '8',         # Use more threads
                '-preset', 'ultrafast',  # Fastest encoding preset
                '-loglevel', 'error',    # Reduce logging for speed
                '-nostats',              # Disable stats for speed
                output_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode != 0:
                logger.error(f"FFmpeg conversion failed: {result.stderr}")
                return audio_path
            
            logger.info(f"Successfully converted to {output_path}")
            return output_path
            
        except subprocess.TimeoutExpired:
            logger.error("FFmpeg conversion timed out")
            return audio_path
        except Exception as e:
            logger.error(f"Error converting audio format: {str(e)}")
            return audio_path
    
    def preprocess_audio_for_speed(self, audio_path: str) -> str:
        """Preprocess audio to speed up transcription with aggressive optimization."""
        try:
            import librosa
            import soundfile as sf
            
            # Load audio with even lower sample rate for speed
            audio, sr = librosa.load(audio_path, sr=8000)  # 8kHz is sufficient for speech
            
            # Apply aggressive noise reduction and normalization
            # Simple normalization
            audio = librosa.util.normalize(audio)
            
            # Optional: Apply simple noise reduction
            # This is a basic noise gate - you can make it more sophisticated
            noise_threshold = 0.01
            audio[abs(audio) < noise_threshold] = 0
            
            # Save preprocessed audio
            preprocessed_path = audio_path.replace('.', '_preprocessed.')
            sf.write(preprocessed_path, audio, sr)
            
            logger.info(f"Audio preprocessed and saved to: {preprocessed_path}")
            return preprocessed_path
            
        except Exception as e:
            logger.warning(f"Audio preprocessing failed: {str(e)}")
            return audio_path  # Return original if preprocessing fails
    
    def split_audio_to_chunks(self, audio_path, chunk_length_ms=600000):
        """Split audio into chunks of chunk_length_ms (default 10 min) and return list of temp file paths."""
        audio = AudioSegment.from_file(audio_path)
        chunks = []
        for i in range(0, len(audio), chunk_length_ms):
            chunk = audio[i:i+chunk_length_ms]
            temp_fd, temp_path = tempfile.mkstemp(suffix='.wav')
            os.close(temp_fd)
            chunk.export(temp_path, format="wav")
            chunks.append(temp_path)
        return chunks

    def transcribe_audio(self, audio_path: str, method: str = "whisper") -> str:
        """Transcribe audio with chunking for long files (Streamlit safe)."""
        try:
            import librosa
            duration = librosa.get_duration(path=audio_path)
            if duration > 600:  # If longer than 10 minutes, chunk
                logger.info(f"Audio longer than 10 minutes ({duration/60:.2f} min), chunking...")
                chunk_paths = self.split_audio_to_chunks(audio_path, chunk_length_ms=600000)
                full_transcript = ""
                for idx, chunk_path in enumerate(chunk_paths):
                    logger.info(f"Transcribing chunk {idx+1}/{len(chunk_paths)}: {chunk_path}")
                    # Preprocess for speed
                    preprocessed_path = self.preprocess_audio_for_speed(chunk_path)
                    transcript = self.transcribe_with_whisper(preprocessed_path)
                    full_transcript += transcript + "\n"
                    # Clean up temp files
                    for temp in [chunk_path, preprocessed_path]:
                        if temp and os.path.exists(temp):
                            try:
                                os.remove(temp)
                            except Exception as e:
                                logger.warning(f"Could not remove temp file {temp}: {e}")
                return full_transcript.strip()
            else:
                # Original logic for short files
                if method == "whisper":
                    if self.ffmpeg_available:
                        logger.info(f"Converting audio to clean WAV format: {audio_path}")
                        clean_wav_path = self.convert_audio_format(audio_path, "wav")
                        if clean_wav_path != audio_path:
                            logger.info(f"Successfully converted to: {clean_wav_path}")
                            audio_path = clean_wav_path
                        else:
                            logger.warning("FFmpeg conversion failed, using original file")
                    else:
                        logger.warning("FFmpeg not available, using original file")
                    preprocessed_path = self.preprocess_audio_for_speed(audio_path)
                    transcription = self.transcribe_with_whisper(preprocessed_path)
                    for temp_path in [preprocessed_path, clean_wav_path if 'clean_wav_path' in locals() else None]:
                        if temp_path and temp_path != audio_path and os.path.exists(temp_path):
                            try:
                                os.remove(temp_path)
                            except:
                                pass
                    return transcription
                else:
                    return self.transcribe_with_speechrecognition(audio_path)
        except Exception as e:
            logger.error(f"Error in transcribe_audio: {str(e)}")
            return f"Error transcribing audio: {str(e)}"
    
    def transcribe_with_speechrecognition(self, audio_path: str) -> str:
        """Transcribe audio using SpeechRecognition (Google Speech API)."""
        try:
            logger.info(f"Transcribing audio with SpeechRecognition: {audio_path}")
            
            with sr.AudioFile(audio_path) as source:
                audio = self.recognizer.record(source)
                transcription = self.recognizer.recognize_google(audio)
            
            logger.info(f"SpeechRecognition transcription completed: {len(transcription)} characters")
            return transcription
            
        except Exception as e:
            logger.error(f"Error transcribing with SpeechRecognition: {str(e)}")
            return f"Error transcribing audio: {str(e)}"
    
    def create_transcript_file(self, transcription: str, original_filename: str) -> str:
        """Create a text file from transcription."""
        try:
            # Create temp directory
            temp_dir = "temp"
            os.makedirs(temp_dir, exist_ok=True)
            
            # Create filename
            base_name = os.path.splitext(original_filename)[0]
            transcript_path = os.path.join(temp_dir, f"{base_name}_transcript.txt")
            
            # Write transcription to file
            with open(transcript_path, 'w', encoding='utf-8') as f:
                f.write(transcription)
            
            logger.info(f"Transcript saved to: {transcript_path}")
            return transcript_path
            
        except Exception as e:
            logger.error(f"Error creating transcript file: {str(e)}")
            return None

    def _create_transcript_pdf(self, transcription: str, pdf_path: str, filename: str):
        """Create a PDF file from transcription using reportlab."""
        try:
            # Create the PDF document
            doc = SimpleDocTemplate(pdf_path, pagesize=letter)
            story = []
            
            # Get styles
            styles = getSampleStyleSheet()
            
            # Create custom styles
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=styles['Heading1'],
                fontSize=16,
                spaceAfter=30,
                alignment=TA_CENTER,
                textColor=blue
            )
            
            header_style = ParagraphStyle(
                'CustomHeader',
                parent=styles['Heading2'],
                fontSize=12,
                spaceAfter=20,
                textColor=black
            )
            
            body_style = ParagraphStyle(
                'CustomBody',
                parent=styles['Normal'],
                fontSize=11,
                spaceAfter=12,
                alignment=TA_LEFT,
                textColor=black
            )
            
            # Add title
            title = Paragraph(f"Audio Transcription: {filename}", title_style)
            story.append(title)
            story.append(Spacer(1, 20))
            
            # Add metadata
            metadata = f"""
            <b>File:</b> {filename}<br/>
            <b>Transcription Date:</b> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}<br/>
            <b>Character Count:</b> {len(transcription):,}<br/>
            <b>Word Count:</b> {len(transcription.split()):,}
            """
            meta_para = Paragraph(metadata, header_style)
            story.append(meta_para)
            story.append(Spacer(1, 30))
            
            # Add transcription content
            # Split transcription into paragraphs for better formatting
            paragraphs = transcription.split('\n\n')
            
            for para in paragraphs:
                if para.strip():
                    # Clean up the paragraph
                    clean_para = para.strip().replace('\n', ' ')
                    if clean_para:
                        p = Paragraph(clean_para, body_style)
                        story.append(p)
                        story.append(Spacer(1, 12))
            
            # Build the PDF
            doc.build(story)
            logger.info(f"PDF transcript created: {pdf_path}")
            
        except Exception as e:
            logger.error(f"Error creating PDF transcript: {str(e)}")
            # If PDF creation fails, just continue with text file
            pass

    def _transcribe_long_audio_chunked(self, audio_path: str, duration: float) -> str:
        """Transcribe long audio files in parallel chunks."""
        try:
            import librosa
            import soundfile as sf
            from concurrent.futures import ThreadPoolExecutor, as_completed
            
            # Load audio
            audio, sr = librosa.load(audio_path, sr=16000)
            
            # Split into 5-minute chunks
            chunk_duration = 300  # 5 minutes in seconds
            chunk_samples = int(chunk_duration * sr)
            
            # Prepare chunks
            chunks = []
            for i in range(0, len(audio), chunk_samples):
                chunk = audio[i:i + chunk_samples]
                chunk_start = i / sr
                chunk_end = min((i + chunk_samples) / sr, duration)
                chunk_num = len(chunks) + 1
                
                # Save chunk to temporary file
                chunk_path = f"temp_chunk_{chunk_num}.wav"
                sf.write(chunk_path, chunk, sr, subtype='PCM_16')
                chunks.append((chunk_num, chunk_path, chunk_start, chunk_end))
            
            print(f" Processing {len(chunks)} chunks in parallel...")
            
            # Process chunks in parallel
            transcriptions = [None] * len(chunks)
            
            def transcribe_chunk(chunk_info):
                chunk_num, chunk_path, chunk_start, chunk_end = chunk_info
                try:
                    print(f" Processing chunk {chunk_num}: {chunk_start/60:.1f}m - {chunk_end/60:.1f}m")
                    
                    result = self.whisper_model.transcribe(
                        chunk_path,
                        language="en",
                        task="transcribe",
                        fp16=False
                    )
                    
                    chunk_text = result.get('text', '').strip()
                    
                    # Clean up temporary file
                    os.remove(chunk_path)
                    
                    if chunk_text:
                        print(f"✅ Chunk {chunk_num} completed: {len(chunk_text)} characters")
                        return chunk_num - 1, chunk_text
                    else:
                        print(f"⚠️ Chunk {chunk_num} returned no text")
                        return chunk_num - 1, ""
                    
                except Exception as e:
                    print(f"❌ Error transcribing chunk {chunk_num}: {str(e)}")
                    if os.path.exists(chunk_path):
                        os.remove(chunk_path)
                    return chunk_num - 1, ""
            
            # Use ThreadPoolExecutor for parallel processing
            with ThreadPoolExecutor(max_workers=4) as executor:
                future_to_chunk = {executor.submit(transcribe_chunk, chunk): chunk for chunk in chunks}
                
                for future in as_completed(future_to_chunk):
                    chunk_idx, text = future.result()
                    if text:
                        transcriptions[chunk_idx] = text
            
            # Combine transcriptions in order
            final_transcription = " ".join([t for t in transcriptions if t])
            
            if final_transcription:
                print(f" Full transcription completed! Total: {len(final_transcription)} characters")
                return final_transcription
            else:
                print("❌ No transcription completed")
                return "No speech detected in audio file."
            
        except Exception as e:
            logger.error(f"Error in parallel chunked transcription: {str(e)}")
            return f"Error transcribing long audio: {str(e)}"

    def transcribe_long_audio_parallel(self, audio_path: str, chunk_duration: int = 300) -> str:
        """Transcribe long audio files using parallel processing for speed."""
        try:
            import librosa
            import soundfile as sf
            from concurrent.futures import ThreadPoolExecutor, as_completed
            
            # Load audio
            audio, sr = librosa.load(audio_path, sr=16000)
            duration = len(audio) / sr
            
            logger.info(f"Processing {duration/60:.1f} minute audio in parallel chunks")
            
            # Split into chunks (5 minutes each)
            chunk_samples = int(chunk_duration * sr)
            
            # Prepare chunks
            chunks = []
            for i in range(0, len(audio), chunk_samples):
                chunk = audio[i:i + chunk_samples]
                chunk_start = i / sr
                chunk_end = min((i + chunk_samples) / sr, duration)
                chunk_num = len(chunks) + 1
                
                # Save chunk to temporary file
                chunk_path = f"temp_chunk_{chunk_num}.wav"
                sf.write(chunk_path, chunk, sr, subtype='PCM_16')
                chunks.append((chunk_num, chunk_path, chunk_start, chunk_end))
            
            print(f" Processing {len(chunks)} chunks in parallel...")
            
            # Process chunks in parallel
            transcriptions = [None] * len(chunks)
            
            def transcribe_chunk(chunk_info):
                chunk_num, chunk_path, chunk_start, chunk_end = chunk_info
                try:
                    print(f" Processing chunk {chunk_num}: {chunk_start/60:.1f}m - {chunk_end/60:.1f}m")
                    
                    # Use tiny model for speed
                    if self.whisper_model is None:
                        self.load_whisper_model("tiny")
                    
                    result = self.whisper_model.transcribe(
                        chunk_path,
                        language="en",
                        task="transcribe",
                        verbose=False,
                        fp16=False,
                        condition_on_previous_text=False,
                        temperature=0.0,
                        compression_ratio_threshold=1.0,
                        logprob_threshold=-2.0,
                        no_speech_threshold=0.8
                    )
                    
                    chunk_text = result.get('text', '').strip()
                    
                    # Clean up temporary file
                    os.remove(chunk_path)
                    
                    if chunk_text:
                        print(f"✅ Chunk {chunk_num} completed: {len(chunk_text)} characters")
                        return chunk_num - 1, chunk_text
                    else:
                        print(f"⚠️ Chunk {chunk_num} returned no text")
                        return chunk_num - 1, ""
                    
                except Exception as e:
                    print(f"❌ Error transcribing chunk {chunk_num}: {str(e)}")
                    if os.path.exists(chunk_path):
                        os.remove(chunk_path)
                    return chunk_num - 1, ""
            
            # Use ThreadPoolExecutor for parallel processing
            with ThreadPoolExecutor(max_workers=4) as executor:
                future_to_chunk = {executor.submit(transcribe_chunk, chunk): chunk for chunk in chunks}
                
                for future in as_completed(future_to_chunk):
                    chunk_idx, text = future.result()
                    if text:
                        transcriptions[chunk_idx] = text
            
            # Combine transcriptions in order
            final_transcription = " ".join([t for t in transcriptions if t])
            
            logger.info(f"Parallel transcription completed: {len(final_transcription)} characters")
            return final_transcription
            
        except Exception as e:
            logger.error(f"Error in parallel transcription: {str(e)}")
            return f"Error in parallel transcription: {str(e)}"

class WebFileHandler:
    def __init__(self):
        self.text_processor = TextProcessor()
        self.audio_processor = AudioProcessor()
        self.saved_pdf_paths = []
        self.processing_status = {}
        self.is_processing = False

    def _process_single_file(self, uploaded_file):
        # Prevent double processing
        if self.is_processing:
            logger.info(f"Already processing a file, skipping {uploaded_file.name}")
            return None
            
        self.is_processing = True
        
        try:
            # Check if this audio file has already been processed
            if uploaded_file.name in self.processing_status:
                logger.info(f"Audio file {uploaded_file.name} already processed, skipping")
                self.is_processing = False
                return None
                
            logger.info(f"Processing file: {uploaded_file.name}")
            ext = os.path.splitext(uploaded_file.name)[1].lower()
            temp_path = f"temp/{uploaded_file.name}"
            os.makedirs("temp", exist_ok=True)
            
            with open(temp_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            # Handle different file types
            if ext in ['.pdf', '.docx', '.pptx', '.xlsx', '.xls']:
                # Existing document processing
                if ext == ".pdf":
                    text_content = self.text_processor.extract_text_from_pdf(temp_path, enable_image_processing=True)
                elif ext == ".docx":
                    text_content = self.text_processor.extract_text_from_docx(temp_path)
                elif ext in [".pptx"]:
                    text_content = self.text_processor.extract_text_from_pptx(temp_path)
                elif ext in [".xlsx"]:
                    text_content = self.text_processor.extract_text_from_xlsx(temp_path)
                elif ext in [".xls"]:
                    text_content = self.text_processor.extract_text_from_xls(temp_path)
                
                self.saved_pdf_paths.append(temp_path)
                
            elif ext in ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.aac']:
                # Audio file processing with speed optimization
                logger.info(f"Processing audio file: {uploaded_file.name}")
                
                # Check file size to decide on processing method
                file_size = os.path.getsize(temp_path)
                file_size_mb = file_size / (1024 * 1024)
                
                if file_size_mb > 50:  # Large files (>50MB)
                    logger.info("Large audio file detected, using parallel processing")
                    transcription = self.audio_processor.transcribe_long_audio_parallel(temp_path)
                else:
                    # Smaller files use regular processing
                    transcription = self.audio_processor.transcribe_audio(temp_path, method="whisper")
                
                logger.info(f"Transcription result: {len(transcription)} characters")
                
                # Check if transcription failed
                if transcription.startswith("Error"):
                    logger.error(f"Transcription failed: {transcription}")
                    self.is_processing = False
                    return None
                
                # Create transcript file
                transcript_path = self.audio_processor.create_transcript_file(transcription, uploaded_file.name)
                if transcript_path:
                    self.saved_pdf_paths.append(transcript_path)
                
                # Process transcription as text content
                text_content = transcription
                self.processing_status[uploaded_file.name] = {
                    'status': 'completed',
                    'transcript_path': transcript_path,
                    'characters': len(transcription)
                }

            else:
                logger.warning(f"Unsupported file type: {uploaded_file.name}")
                self.is_processing = False
                return None

            if text_content and text_content.strip():
                logger.info(f"Processing text content: {len(text_content)} characters")
                chunks = self.text_processor.chunk_text(text_content)
                documents = []
                for i, chunk in enumerate(chunks):
                    if chunk.strip():
                        documents.append({
                            'text': chunk.strip(),
                            'metadata': {
                                'source': uploaded_file.name,
                                'chunk_id': i,
                                'total_chunks': len(chunks),
                                'date': datetime.now().strftime('%Y-%m-%d'),
                                'file_type': 'audio' if ext in ['.mp3', '.wav', '.m4a', '.flac', '.ogg', '.aac'] else 'document'
                            }
                        })
                logger.info(f"Successfully processed {uploaded_file.name}: {len(documents)} chunks")
                self.is_processing = False
                return documents
            else:
                logger.warning(f"No text content extracted from {uploaded_file.name}")
                self.is_processing = False
                return None

        except Exception as e:
            import traceback
            error_details = traceback.format_exc()
            logger.error(f"Error processing file {uploaded_file.name}: {str(e)}")
            logger.error(f"Full traceback: {error_details}")
            self.is_processing = False
            return None

    def get_saved_pdf_paths(self):
        """Get the paths of saved PDF files for image processing."""
        return self.saved_pdf_paths

    def get_processing_status(self):
        """Get the processing status of uploaded files."""
        return self.processing_status

### =================== Vector Store =================== ###
class VectorStore:
    _model = None  # Class variable to store the model
    
    def __init__(self, dimension: int = 768):
        self.documents = []
        self.embeddings = None
        self.index = None
        self.chunk_size = 2000
        self.chunk_overlap = 100
        self.initialized = False
        
        # Initialize embedding model only once
        self._initialize_model()

    def _initialize_model(self):
        """Initialize the embedding model only once."""
        if VectorStore._model is None:  # Only load if not already loaded
            try:
                from sentence_transformers import SentenceTransformer
                import torch
                import os
                import streamlit as st
                
                # Get Hugging Face token from Streamlit secrets
                hf_token = st.secrets.get("HUGGING_FACE_HUB_TOKEN")
                if hf_token:
                    # Set the token for Hugging Face Hub
                    os.environ['HF_TOKEN'] = hf_token
                    # Also set it for the transformers library
                    os.environ['HUGGING_FACE_HUB_TOKEN'] = hf_token
                    logger.info("Hugging Face premium token configured from Streamlit secrets")
                else:
                    logger.warning("No Hugging Face token found in Streamlit secrets - using free tier")
                
                # Force CPU usage and handle meta tensor issue
                logger.info("Loading sentence transformers model...")
                
                # Set environment variables to force CPU usage
                os.environ['CUDA_VISIBLE_DEVICES'] = ''
                
                # Try different approaches to handle the device issue
                try:
                    # Method 1: Load with explicit CPU device and trust_remote_code
                    VectorStore._model = SentenceTransformer(
                        'all-MiniLM-L6-v2', 
                        device='cpu',
                        trust_remote_code=True
                    )
                    logger.info("Method 1 successful: Loaded with CPU device")
                    
                except Exception as e1:
                    logger.warning(f"Method 1 failed: {e1}")
                    try:
                        # Method 2: Load normally then move to CPU
                        VectorStore._model = SentenceTransformer('all-MiniLM-L6-v2')
                        VectorStore._model.to('cpu')
                        logger.info("Method 2 successful: Loaded normally and moved to CPU")
                        
                    except Exception as e2:
                        logger.warning(f"Method 2 failed: {e2}")
                        try:
                            # Method 3: Load with specific torch settings
                            torch.set_default_device('cpu')
                            VectorStore._model = SentenceTransformer('all-MiniLM-L6-v2')
                            VectorStore._model.to('cpu')
                            logger.info("Method 3 successful: Used torch default device")
                            
                        except Exception as e3:
                            logger.warning(f"Method 3 failed: {e3}")
                            # Method 4: Last resort - try with different model
                            try:
                                VectorStore._model = SentenceTransformer('paraphrase-MiniLM-L3-v2')
                                VectorStore._model.to('cpu')
                                logger.info("Method 4 successful: Used alternative model")
                            except Exception as e4:
                                logger.warning(f"Method 4 failed: {e4}")
                                # Method 5: Try with even simpler model
                                try:
                                    VectorStore._model = SentenceTransformer('distiluse-base-multilingual-cased-v2')
                                    VectorStore._model.to('cpu')
                                    logger.info("Method 5 successful: Used multilingual model")
                                except Exception as e5:
                                    logger.error(f"All methods failed. Last error: {e5}")
                                    self.initialized = False
                                    return
                
                # Verify the model is working
                try:
                    test_embedding = VectorStore._model.encode(['test'], convert_to_tensor=False)
                    if test_embedding is not None and len(test_embedding) > 0:
                        logger.info(f"Model test successful. Embedding shape: {test_embedding.shape}")
                    else:
                        logger.error("Model test failed: Empty embedding returned")
                        self.initialized = False
                        return
                except Exception as e:
                    logger.error(f"Model test failed: {e}")
                    self.initialized = False
                    return
                
                logger.info("Hugging Face embedding model initialized successfully")
                
            except Exception as e:
                logger.error(f"Failed to initialize Hugging Face model: {str(e)}")
                self.initialized = False
                return
        
        self.model = VectorStore._model  # Use the shared model
        self.initialized = True

    def get_embeddings_batch(self, texts):
        """Get embeddings from Hugging Face model in batches."""
        try:
            if not self.initialized or self.model is None:
                logger.error("Model not initialized")
                return None

            # Validate input texts
            if not texts or len(texts) == 0:
                logger.error("No texts provided for embedding")
                return None
            
            # Filter out empty texts
            valid_texts = [text for text in texts if text and text.strip()]
            if len(valid_texts) != len(texts):
                logger.warning(f"Filtered out {len(texts) - len(valid_texts)} empty texts")

            if not valid_texts:
                logger.error("No valid texts after filtering")
                return None

            # Get embeddings for all texts at once
            embeddings = self.model.encode(valid_texts, convert_to_tensor=False)
            
            # Validate embeddings
            if embeddings is None or len(embeddings) == 0:
                logger.error("Model returned empty embeddings")
                return None
                
            # Convert to list format
            if hasattr(embeddings, 'tolist'):
                embeddings = embeddings.tolist()
            
            # Ensure we have the right number of embeddings
            if len(embeddings) != len(valid_texts):
                logger.error(f"Embedding count mismatch: {len(embeddings)} embeddings for {len(valid_texts)} texts")
                return None
                
            return embeddings
            
        except Exception as e:
            logger.error(f"Error getting Hugging Face embeddings: {str(e)}")
            return None

    def add_documents(self, documents: List[Dict]):
        """Add documents to the vector store using Hugging Face embeddings."""
        if not documents:
            return
        
        if not self.initialized:
            logger.error("Model not initialized, cannot add documents")
            logger.error("This usually means the sentence transformers model failed to load")
            logger.error("Check the logs above for Hugging Face model initialization errors")
            return
        
        try:
            texts = [doc['text'] for doc in documents]
            
            # Get embeddings for all documents
            embeddings = self.get_embeddings_batch(texts)
            
            if embeddings is None or len(embeddings) != len(documents):
                logger.error("Failed to get embeddings for all documents")
                return
            
            self.documents.extend(documents)
            embeddings = np.array(embeddings)
            
            if self.embeddings is None:
                self.embeddings = embeddings
            else:
                self.embeddings = np.vstack([self.embeddings, embeddings])
            
            embedding_dim = self.embeddings.shape[1]
            self.index = faiss.IndexFlatIP(embedding_dim)
            self.index.add(self.embeddings.astype('float32'))
            
            logger.info(f"Added {len(documents)} documents to vector store (Hugging Face embeddings)")
            
        except Exception as e:
            logger.error(f"Error adding documents to vector store: {str(e)}")

    def search(self, query: str, k: int = 3) -> List[Dict]:
        """Search for similar documents using Hugging Face embeddings."""
        if not self.documents or self.index is None or not self.initialized:
            return []
        
        try:
            # Get query embedding
            query_embedding = self.model.encode([query], convert_to_tensor=False)
            query_embedding = query_embedding.reshape(1, -1)
            
            # Safety check: ensure we have valid embeddings
            if self.embeddings is None or len(self.embeddings) == 0:
                logger.warning("No embeddings available for search")
                return []
            
            scores, indices = self.index.search(query_embedding.astype('float32'), k)
            results = []
            for i, idx in enumerate(indices[0]):
                if idx < len(self.documents):
                    results.append({
                        'text': self.documents[idx]['text'],
                        'metadata': self.documents[idx].get('metadata', {}),
                        'score': float(scores[0][i])
                    })
            return results
        except Exception as e:
            logger.error(f"Error searching vector store: {str(e)}")
            return []

    def is_ready(self) -> bool:
        """Check if the vector store is properly initialized and ready for use."""
        return (
            self.initialized and
            self.documents is not None and 
            len(self.documents) > 0 and 
            self.embeddings is not None and 
            len(self.embeddings) > 0 and 
            self.index is not None
        )

### =================== Claude Handler =================== ###
class ClaudeHandler:
    def __init__(self, system_prompt=None):
        self.api_url = "https://api.openai.com/v1/chat/completions"
        self.headers = {
            "Authorization": f"Bearer {OPENAI_API_KEY}",
            "Content-Type": "application/json"
        }
        # Enhanced prompt that prioritizes specific, detailed information
        self.system_prompt = (
            system_prompt or
            "You are an expert financial consultant and analyst with 20+ years of experience in investment banking, "
            "equity research, and financial modeling. You specialize in deep-dive financial analysis and comprehensive reporting. "
            
            "CRITICAL: When asked about financial data, revenue, projections, or specific figures, ALWAYS lead with the most "
            "specific and detailed numerical information available. Do NOT start with vague summaries or general statements. "
            "If the context contains specific numbers, dates, or financial figures, present those FIRST. Only provide general "
            "commentary AFTER presenting the specific data. "

            "You must provide comprehensive, detailed analysis in paragraph form. Never give one-sentence answers. "
            "Always provide thorough analysis with multiple data points, historical context, year-over-year comparisons, "
            "trend analysis, risk factors, market conditions, competitive analysis, and investment implications. "
            "When projections are available, explain the assumptions and growth drivers. Compare with industry benchmarks, "
            "competitors, or historical performance. Provide insights on what the data means for investors and stakeholders. "
            
            "Structure your responses with clear paragraphs that flow logically: start with an executive summary of key findings, "
            "present detailed financial metrics with exact numbers, provide historical context and trends, include comparative "
            "analysis, discuss implications and outlook, and address potential risks and opportunities. "
            
            "Always cite the document source (filename and page or chunk) for every fact or statement in your answer. "
            "If you use multiple sources, cite each one clearly. Include page numbers when available. "
            
            "Remember: You are a senior financial analyst. Your clients expect comprehensive, detailed analysis that provides "
            "actionable insights. Never provide superficial or vague responses."
        )
        self.response_cache = {}  # Simple cache for repeated queries

    def generate_answer(self, question: str, context: str, normalize_length: bool = True) -> str:
        try:
            # Check cache first
            cache_key = f"{question[:100]}_{hash(context[:500])}"
            if cache_key in self.response_cache:
                return self.response_cache[cache_key]
            
            # Set maximum tokens for both input and output
            max_input_tokens = 128000  # Maximum context window for GPT-4o
            max_output_tokens = 16384   # Maximum response tokens
            
            # Count and limit input tokens
            import tiktoken
            enc = tiktoken.get_encoding("cl100k_base")
            
            # Count input tokens
            input_text = f"Context: {context}\n\nQuestion: {question}"
            input_tokens = len(enc.encode(input_text))
            
            # Truncate context if too long
            if input_tokens > max_input_tokens:
                # Calculate how much context we can keep
                question_tokens = len(enc.encode(question))
                system_tokens = len(enc.encode(self.system_prompt))
                available_tokens = max_input_tokens - question_tokens - system_tokens - 100  # Buffer
                
                if available_tokens > 0:
                    # Truncate context to fit
                    context = enc.decode(enc.encode(context)[:available_tokens])
                else:
                    return "Error: Question too long for processing."
            
            messages = [
                {
                    "role": "system",
                    "content": self.system_prompt
                },
                {
                    "role": "user",
                    "content": f"Context: {context}\n\nQuestion: {question}"
                }
            ]

            payload = {
                "model": "gpt-4o",
                "messages": messages,
                "temperature": 0.3,
                "max_tokens": max_output_tokens  # Maximum output tokens
            }

            response = requests.post(
                self.api_url,
                headers=self.headers,
                json=payload,
                timeout=100
            )

            response.raise_for_status()
            response_data = response.json()
            answer = response_data["choices"][0]["message"]["content"]
            
            # Cache the response
            self.response_cache[cache_key] = answer
            
            return answer

        except Exception as e:
            logger.error(f"Error generating answer: {str(e)}")
            return f"Error generating answer: {str(e)}"

    def _get_response_length(self, question: str, normalize_length: bool) -> int:
        """Determine appropriate response length based on question type."""
        if not normalize_length:
            return 16384
        
        question_lower = question.lower()
        
        # Short responses for simple questions
        if any(word in question_lower for word in ['what is', 'define', 'explain briefly', 'summarize']):
            return 16384
        
        # Medium responses for analysis questions
        if any(word in question_lower for word in ['analyze', 'compare', 'discuss', 'evaluate']):
            return 16384
        
        # Long responses for complex questions
        if any(word in question_lower for word in ['detailed', 'comprehensive', 'thorough']):
            return 16384
        
        # Default medium length
        return 16384

### =================== Question Handler =================== ###
class QuestionHandler:
    def __init__(self, vector_store: VectorStore):
        self.vector_store = vector_store
        self.llm = ClaudeHandler()
        self.conversation_history = []
        self.error_count = 0
        self.max_retries = 3

    def process_question(self, question: str, query_type: str = "document", k: int = 5, normalize_length: bool = True, analysis_mode: str = "General") -> str:
        for attempt in range(self.max_retries):
            try:
                # DETECT FRENCH QUESTION
                is_french = self._is_french_question(question)
                
                results = self.vector_store.search(question, k=k)
                if not results:
                    return "No relevant information found in the documents." if not is_french else "Aucune information pertinente trouvée dans les documents."
                
                # Build context with metadata for citation
                context = "\n".join([
                    f"[source: {chunk['metadata'].get('source', 'unknown')}, chunk: {chunk['metadata'].get('chunk_id', '?')}]\n{chunk['text']}"
                    for chunk in results
                ])
                
                # Choose system prompt based on analysis_mode and language
                if is_french:
                    # French system prompts - COMPLETE AND PROPER
                    if analysis_mode == "Financial Document":
                        system_prompt = (
                            "Vous êtes un analyste financier senior avec plus de 20 ans d'expérience dans les banques d'investissement de premier plan et les cabinets de conseil. "
                            "Votre analyse doit être EXTREMEMENT approfondie et complète - fournissez des réponses détaillées de qualité ChatGPT. "
                            "Rédigez des rapports en paragraphes détaillés, PAS de points ou de résumés brefs. Utilisez des paragraphes fluides et détaillés qui fournissent une analyse complète.\n\n"
                            "IMPORTANT: Répondez UNIQUEMENT en français. Ne mélangez jamais l'anglais et le français dans votre réponse.\n\n"
                            "Pour chaque document financier que vous analysez:\n\n"
                            "1. **Résumé exécutif**: Rédigez 2-3 paragraphes détaillés mettant en évidence les insights financiers les plus critiques avec le contexte complet et les implications\n"
                            "2. **Analyse de la performance financière**: Rédigez des paragraphes complets couvrant:\n"
                            "   Les tendances de revenus, taux de croissance et moteurs avec une analyse détaillée de ce qui stimule la performance\n"
                            "   Les métriques de rentabilité (marge brute, marge opérationnelle, marge nette) avec le contexte complet et les comparaisons sectorielles\n"
                            "   L'analyse des flux de trésorerie (exploitation, investissement, financement) avec une explication détaillée de la génération et de l'utilisation de trésorerie\n"
                            "   La solidité du bilan et la liquidité avec une évaluation complète de la position financière\n"
                            "   Les ratios financiers clés et les benchmarks avec une interprétation détaillée et la signification\n"
                            "3. **Contexte historique**: Rédigez des paragraphes détaillés comparant la performance actuelle aux tendances historiques (3-5 ans) avec une analyse complète des changements et des moteurs\n"
                            "4. **Analyse comparative**: Rédigez des paragraphes complets de benchmarking contre les pairs du secteur et les moyennes du marché avec un positionnement concurrentiel détaillé\n"
                            "5. **Évaluation des risques**: Rédigez des paragraphes détaillés identifiant les risques financiers, de marché, opérationnels avec une explication complète des impacts potentiels\n"
                            "6. **Perspectives futures**: Rédigez des paragraphes complets couvrant les projections, prévisions et déclarations prospectives avec une justification détaillée\n"
                            "7. **Implications d'investissement**: Rédigez des paragraphes détaillés expliquant ce que cela signifie pour les investisseurs, les parties prenantes avec le contexte complet et les recommandations\n"
                            "8. **Insights actionnables**: Rédigez des paragraphes complets avec des recommandations spécifiques et les prochaines étapes avec des conseils d'implémentation détaillés\n\n"
                            "Citez toujours la source du document (nom de fichier et page/segment) pour chaque fait, chiffre ou déclaration. "
                            "Utilisez des chiffres spécifiques, des pourcentages et des dates. Fournissez le contexte complet pour toutes les métriques. "
                            "Rédigez en paragraphes fluides et détaillés qui se lisent comme un rapport financier professionnel. "
                            "Visez des réponses de 800-1500 mots minimum, avec une couverture complète de tous les aspects financiers pertinents en format paragraphe."
                        )
                    elif analysis_mode == "Company Evaluation":
                        system_prompt = (
                            "Vous êtes un consultant d'entreprise senior et stratège corporatif avec une expertise en évaluation d'entreprise, analyse concurrentielle et planification stratégique. "
                            "Votre évaluation doit être EXTREMEMENT approfondie et complète - fournissez des réponses détaillées de qualité ChatGPT. "
                            "Rédigez des rapports en paragraphes détaillés, PAS de points ou de résumés brefs. Utilisez des paragraphes fluides et détaillés qui fournissent une analyse complète.\n\n"
                            "IMPORTANT: Répondez UNIQUEMENT en français. Ne mélangez jamais l'anglais et le français dans votre réponse.\n\n"
                            "Pour chaque évaluation d'entreprise:\n\n"
                            "1. **Aperçu de l'entreprise**: Rédigez des paragraphes complets couvrant le profil de l'entreprise, le modèle d'affaires et la position de marché avec le contexte stratégique complet\n"
                            "2. **Analyse SWOT**: Rédigez des paragraphes détaillés analysant:\n"
                            "   Forces: Avantages concurrentiels, capacités uniques, leadership de marché avec une explication complète de la valeur stratégique\n"
                            "   Faiblesses: Défis opérationnels, contraintes de ressources, écarts concurrentiels avec une analyse d'impact détaillée\n"
                            "   Opportunités: Expansion de marché, nouveaux produits, partenariats stratégiques avec une évaluation complète du potentiel de croissance\n"
                            "   Menaces: Pressions concurrentielles, disruption de marché, risques réglementaires avec une analyse d'impact des risques détaillée\n"
                            "3. **Analyse de marché**: Rédigez des paragraphes complets couvrant la dynamique sectorielle, le paysage concurrentiel, l'analyse de part de marché avec le contexte de marché complet\n"
                            "4. **Évaluation du leadership**: Rédigez des paragraphes détaillés évaluant la qualité de l'équipe de direction, la vision stratégique, la capacité d'exécution avec une analyse de leadership complète\n"
                            "5. **Santé financière**: Rédigez des paragraphes complets couvrant la stabilité des revenus, la rentabilité, les flux de trésorerie, les niveaux de dette avec une analyse financière détaillée\n"
                            "6. **Perspectives de croissance**: Rédigez des paragraphes détaillés analysant les opportunités de croissance organique, le potentiel d'expansion, le pipeline d'innovation avec une évaluation complète de la stratégie de croissance\n"
                            "7. **Facteurs de risque**: Rédigez des paragraphes complets couvrant les risques opérationnels, financiers, de marché et stratégiques avec une analyse des risques détaillée\n"
                            "8. **Recommandations stratégiques**: Rédigez des paragraphes détaillés avec des actions spécifiques pour l'amélioration et la croissance avec des conseils d'implémentation complets\n"
                            "9. **Perspective d'évaluation**: Rédigez des paragraphes complets couvrant l'évaluation de la juste valeur, l'attractivité d'investissement avec une analyse d'évaluation détaillée\n"
                            "10. **Perspectives futures**: Rédigez des paragraphes complets couvrant les projections de 3-5 ans et la trajectoire stratégique avec une analyse prospective complète\n\n"
                            "Citez toujours la source du document (nom de fichier et page/segment) pour chaque fait ou déclaration. "
                            "Utilisez des exemples spécifiques, des points de données et des benchmarks sectoriels. "
                            "Fournissez des insights actionnables et des recommandations stratégiques en format paragraphe détaillé. "
                            "Rédigez en paragraphes fluides et complets qui se lisent comme un rapport d'évaluation d'entreprise professionnel. "
                            "Visez des réponses de 1000-1800 mots minimum, couvrant tous les aspects de l'entreprise de manière complète en format paragraphe."
                        )
                    else:
                        system_prompt = (
                            "Vous êtes un analyste d'entreprise hautement qualifié et consultant avec une expertise en finance, technologie et stratégie. "
                            "Votre tâche est de rédiger un rapport d'entreprise long format, profondément analytique et riche en insights (minimum 1 000 mots) "
                            "basé sur le contenu du document fourni. Votre réponse doit se lire comme un mémo d'intelligence d'entreprise de niveau professionnel "
                            "de premier plan—clair, structuré et convaincant, mais jamais superficiel ou artificiel.\n\n"
                            "IMPORTANT: Répondez UNIQUEMENT en français. Ne mélangez jamais l'anglais et le français dans votre réponse.\n\n"
                            "DIRECTIVES IMPORTANTES:\n\n"
                            "Utilisez une structure narrative fluide en paragraphes (pas de points ou d'en-têtes de section), mais avec des transitions logiques entre les thèmes.\n\n"
                            "Citez le nom de fichier du document source et le numéro de segment/page pour chaque affirmation ou insight (ex: \"source: Rue de Vernaz.m4a, segment 16\").\n\n"
                            "Ne vous contentez pas de résumer. Analysez plutôt les implications, décomposez le contexte stratégique, explorez les risques et opportunités potentiels, "
                            "et déduisez la logique non énoncée derrière les déclarations clés.\n\n"
                            "Priorisez la spécificité par rapport aux généralités. Citez ou paraphrasez les idées importantes et expliquez pourquoi elles comptent.\n\n"
                            "Reliez les observations en une narration cohérente. Ne listez pas simplement les idées—connectez-les en une histoire stratégique qui reflète "
                            "une compréhension profonde de l'entreprise, du paysage concurrentiel et des tendances évolutives.\n\n"
                            "Soyez critique, créatif et réaliste. Mettez en évidence les tensions, risques ou compromis, et offrez des insights nuancés.\n\n"
                            "Absolument aucune phrase de remplissage vague ou robotique comme \"Voici un rapport détaillé...\" ou \"Comme vous l'avez demandé...\"\n\n"
                            "Votre ton doit refléter la maîtrise, la clarté et l'insight—comme un consultant McKinsey écrivant pour un CIO de fonds spéculatif. "
                            "Assurez-vous que le résultat est une analyse d'entreprise autonome, hautement lisible qui délivre une vraie valeur stratégique."
                        )
                else:
                    # English system prompts (existing code)
                    if analysis_mode == "Financial Document":
                        system_prompt = (
                            "You are a senior financial analyst with 20+ years of experience at top-tier investment banks and consulting firms. "
                            "Your analysis must be EXTREMELY thorough and comprehensive - provide ChatGPT-quality detailed responses. "
                            "Write in-depth paragraph reports, NOT bullet points or brief summaries. Use flowing, detailed paragraphs that provide comprehensive analysis.\n\n"
                            "For every financial document you analyze:\n\n"
                            "1. **Executive Summary**: Write 2-3 detailed paragraphs highlighting the most critical financial insights with full context and implications\n"
                            "2. **Financial Performance Analysis**: Write comprehensive paragraphs covering:\n"
                            "   Revenue trends, growth rates, and drivers with detailed analysis of what's driving performance\n"
                            "   Profitability metrics (gross margin, operating margin, net margin) with full context and industry comparisons\n"
                            "   Cash flow analysis (operating, investing, financing) with detailed explanation of cash generation and usage\n"
                            "   Balance sheet strength and liquidity with comprehensive assessment of financial position\n"
                            "   Key financial ratios and benchmarks with detailed interpretation and significance\n"
                            "3. **Historical Context**: Write detailed paragraphs comparing current performance to historical trends (3-5 years) with full analysis of changes and drivers\n"
                            "4. **Comparative Analysis**: Write comprehensive paragraphs benchmarking against industry peers and market averages with detailed competitive positioning\n"
                            "5. **Risk Assessment**: Write detailed paragraphs identifying financial risks, market risks, operational risks with full explanation of potential impacts\n"
                            "6. **Future Outlook**: Write comprehensive paragraphs covering projections, forecasts, and forward-looking statements with detailed rationale\n"
                            "7. **Investment Implications**: Write detailed paragraphs explaining what this means for investors, stakeholders with full context and recommendations\n"
                            "8. **Actionable Insights**: Write comprehensive paragraphs with specific recommendations and next steps with detailed implementation guidance\n\n"
                            "Always cite the document source (filename and page/chunk) for every fact, figure, or statement. "
                            "Use specific numbers, percentages, and dates. Provide full context for all metrics. "
                            "Write in flowing, detailed paragraphs that read like a professional financial report. "
                            "Aim for responses that are 800-1500 words minimum, with comprehensive coverage of all relevant financial aspects in paragraph format."
                        )
                    elif analysis_mode == "Company Evaluation":
                        system_prompt = (
                            "You are a senior business consultant and corporate strategist with expertise in company valuation, competitive analysis, and strategic planning. "
                            "Your evaluation must be EXTREMELY thorough and comprehensive - provide ChatGPT-quality detailed responses. "
                            "Write in-depth paragraph reports, NOT bullet points or brief summaries. Use flowing, detailed paragraphs that provide comprehensive analysis.\n\n"
                            "For every company evaluation:\n\n"
                            "1. **Company Overview**: Write comprehensive paragraphs covering company profile, business model, and market position with full strategic context\n"
                            "2. **SWOT Analysis**: Write detailed paragraphs analyzing:\n"
                            "   Strengths: Competitive advantages, unique capabilities, market leadership with full explanation of strategic value\n"
                            "   Weaknesses: Operational challenges, resource constraints, competitive gaps with detailed impact analysis\n"
                            "   Opportunities: Market expansion, new products, strategic partnerships with comprehensive growth potential assessment\n"
                            "   Threats: Competitive pressures, market disruption, regulatory risks with detailed risk impact analysis\n"
                            "3. **Market Analysis**: Write comprehensive paragraphs covering industry dynamics, competitive landscape, market share analysis with full market context\n"
                            "4. **Leadership Assessment**: Write detailed paragraphs evaluating management team quality, strategic vision, execution capability with full leadership analysis\n"
                            "5. **Financial Health**: Write comprehensive paragraphs covering revenue stability, profitability, cash flow, debt levels with detailed financial analysis\n"
                            "6. **Growth Prospects**: Write detailed paragraphs analyzing organic growth opportunities, expansion potential, innovation pipeline with comprehensive growth strategy assessment\n"
                            "7. **Risk Factors**: Write comprehensive paragraphs covering operational, financial, market, and strategic risks with detailed risk analysis\n"
                            "8. **Strategic Recommendations**: Write detailed paragraphs with specific actions for improvement and growth with comprehensive implementation guidance\n"
                            "9. **Valuation Perspective**: Write comprehensive paragraphs covering fair value assessment, investment attractiveness with detailed valuation analysis\n"
                            "10. **Future Outlook**: Write detailed paragraphs covering 3-5 year projections and strategic trajectory with comprehensive forward-looking analysis\n\n"
                            "Always cite the document source (filename and page/chunk) for every fact or statement. "
                            "Use specific examples, data points, and industry benchmarks. "
                            "Provide actionable insights and strategic recommendations in detailed paragraph format. "
                            "Write in flowing, comprehensive paragraphs that read like a professional business evaluation report. "
                            "Aim for responses that are 1000-1800 words minimum, covering all aspects of the business comprehensively in paragraph format."
                        )
                    elif analysis_mode == "Legal Document":
                        system_prompt = (
                            "You are a senior legal expert and attorney with 20+ years of experience in corporate law, regulatory compliance, and legal risk assessment. "
                            "Your legal analysis must be EXTREMELY thorough and comprehensive - provide ChatGPT-quality detailed responses. "
                            "Write in-depth paragraph reports, NOT bullet points or brief summaries. Use flowing, detailed paragraphs that provide comprehensive analysis.\n\n"
                            "For every legal document you analyze:\n\n"
                            "1. **Document Overview**: Write comprehensive paragraphs covering type of legal document, parties involved, effective dates with full legal context\n"
                            "2. **Key Legal Provisions**: Write detailed paragraphs analyzing all important clauses and terms with comprehensive legal interpretation\n"
                            "3. **Obligations and Rights**: Write comprehensive paragraphs clearly breaking down what each party must do and is entitled to with full legal implications\n"
                            "4. **Risk Assessment**: Write detailed paragraphs covering legal risks, potential liabilities, compliance issues with comprehensive risk analysis\n"
                            "5. **Regulatory Compliance**: Write comprehensive paragraphs covering applicable laws, regulations, and compliance requirements with detailed compliance analysis\n"
                            "6. **Enforcement Mechanisms**: Write detailed paragraphs explaining how the agreement is enforced, dispute resolution procedures with comprehensive enforcement analysis\n"
                            "7. **Termination Clauses**: Write detailed paragraphs covering conditions for ending the agreement, notice requirements with detailed termination analysis\n"
                            "8. **Amendments and Modifications**: Write detailed paragraphs explaining how changes can be made to the agreement with comprehensive amendment analysis\n"
                            "9. **Legal Precedents**: Write comprehensive paragraphs covering relevant case law or legal principles that apply with detailed precedent analysis\n"
                            "10. **Practical Implications**: Write detailed paragraphs covering real-world impact on business operations with comprehensive practical analysis\n"
                            "11. **Recommendations**: Write comprehensive paragraphs with suggested actions, areas of concern, negotiation points with detailed guidance\n"
                            "12. **Compliance Checklist**: Write detailed paragraphs covering specific steps needed to ensure legal compliance with comprehensive compliance guidance\n\n"
                            "Always cite the document source (filename and page/chunk) for every legal provision or statement. "
                            "Use precise legal terminology and explain complex concepts clearly in detailed paragraphs. "
                            "Highlight potential legal issues and provide practical guidance in comprehensive paragraph format. "
                            "Write in flowing, detailed paragraphs that read like a professional legal memorandum. "
                            "Aim for responses that are 800-1500 words minimum, with comprehensive legal analysis in paragraph format."
                        )
                    elif analysis_mode == "Financial Excel Document":
                        system_prompt = (
                            "You are a senior financial data analyst and Excel expert with deep expertise in financial modeling, data analysis, and spreadsheet interpretation. "
                            "Your Excel analysis must be EXTREMELY thorough and comprehensive - provide ChatGPT-quality detailed responses. "
                            "Write in-depth paragraph reports, NOT bullet points or brief summaries. Use flowing, detailed paragraphs that provide comprehensive analysis.\n\n"
                            "For every Excel document you analyze:\n\n"
                            "1. **Document Structure**: Write comprehensive paragraphs covering overview of worksheets, data organization, and key tables with full structural analysis\n"
                            "2. **Data Summary**: Write detailed paragraphs providing comprehensive summary of all financial data and key metrics with full data interpretation\n"
                            "3. **Financial Performance Analysis**: Write comprehensive paragraphs covering:\n"
                            "   Revenue analysis by period, segment, or product with detailed performance drivers and trends\n"
                            "   Cost structure and profitability breakdown with comprehensive cost analysis and margin drivers\n"
                            "   Cash flow analysis and working capital trends with detailed cash flow interpretation\n"
                            "   Balance sheet analysis and financial position with comprehensive balance sheet assessment\n"
                            "4. **Trend Analysis**: Write detailed paragraphs identifying patterns, growth rates, and anomalies in the data with comprehensive trend interpretation\n"
                            "5. **Key Metrics Calculation**: Write comprehensive paragraphs covering important ratios, KPIs, and performance indicators with detailed metric analysis\n"
                            "6. **Data Quality Assessment**: Write detailed paragraphs covering accuracy, completeness, and reliability of the data with comprehensive quality analysis\n"
                            "7. **Variance Analysis**: Write comprehensive paragraphs comparing actual vs. budget, period-over-period changes with detailed variance interpretation\n"
                            "8. **Forecasting Insights**: Write detailed paragraphs covering projections, assumptions, and future outlook with comprehensive forecasting analysis\n"
                            "9. **Risk Indicators**: Write detailed paragraphs covering financial stress points, concerning trends, red flags with detailed risk analysis\n"
                            "10. **Actionable Insights**: Write detailed paragraphs with specific recommendations based on the data with comprehensive implementation guidance\n"
                            "11. **Data Visualization Suggestions**: Write comprehensive paragraphs covering charts and graphs that would enhance understanding with detailed visualization analysis\n"
                            "12. **Model Validation**: Write detailed paragraphs covering assessment of formulas, calculations, and model integrity with comprehensive validation analysis\n\n"
                            "Always cite the document source (filename and page/chunk) for every data point or calculation. "
                            "Use specific numbers, percentages, and formulas. Explain the meaning behind the data in detailed paragraphs. "
                            "Highlight important trends and provide context for all metrics in comprehensive paragraph format. "
                            "Write in flowing, detailed paragraphs that read like a professional financial data analysis report. "
                            "Aim for responses that are 1000-1800 words minimum, with comprehensive data analysis in paragraph format."
                        )
                    else:
                        system_prompt = (
                            "You are a highly skilled business analyst and consultant with expertise in finance, technology, and strategy. Your task is to write a long-form, deeply analytical, and insight-rich business report (minimum 1,000 words) based on the contents of the provided document. Your response should read like a top-tier professional business intelligence memo—clear, structured, and compelling, but never fluffy or superficial.\n\n"
                            "IMPORTANT GUIDELINES:\n\n"
                            "Use a flowing narrative paragraph structure (no bullet points or section headers), but with logical transitions between themes.\n\n"
                            "Cite the source document filename and chunk/page number for each assertion or insight (e.g., \"source: Rue de Vernaz.m4a, chunk 16\").\n\n"
                            "Do not merely summarize. Instead, analyze the implications, unpack the strategic context, explore potential risks and opportunities, and infer the unstated logic behind key statements.\n\n"
                            "Prioritize specificity over generalities. Quote or paraphrase important ideas and explain why they matter.\n\n"
                            "Tie observations together into a cohesive narrative. Don't just list ideas—connect them into a strategic storyline that reflects a deep understanding of the business, competitive landscape, and evolving trends.\n\n"
                            "Be critical, creative, and realistic. Highlight tensions, risks, or tradeoffs, and offer nuanced insight.\n\n"
                            "Absolutely no vague or robotic filler phrases like \"Here is a detailed report...\" or \"As you requested...\"\n\n"
                            "Your tone should reflect mastery, clarity, and insight—like a McKinsey consultant writing for a hedge fund CIO. Make sure the result is a standalone, highly readable business analysis that delivers real strategic value."
                        )
                
                enhanced_question = f"{system_prompt}\n\nUser question: {question}\n"
                answer = self.llm.generate_answer(enhanced_question, context, normalize_length=False)
                
                # REMOVE TRANSLATION - Let the AI generate French directly
                # if is_french:
                #     answer = self._translate_to_french(answer)
                
                top_chunk = results[0] if results else None
                self.conversation_history.append({
                    'question': question,
                    'answer': answer,
                    'question_type': 'document',
                    'mode': 'document',
                    'timestamp': datetime.now().isoformat(),
                    'source': top_chunk['metadata'].get('source') if top_chunk else None,
                    'page': top_chunk['metadata'].get('page') if top_chunk else None,
                    'chunk_text': top_chunk['text'] if top_chunk else None
                })
                return answer
            except Exception as e:
                self.error_count += 1
                if attempt == self.max_retries - 1:
                    return f"Error processing question: {str(e)}" if not is_french else f"Erreur lors du traitement de la question: {str(e)}"
                continue

    def process_follow_up(self, follow_up_question: str, k: int = 5, normalize_length: bool = True) -> str:
        """Process a follow-up question using conversation history and document context."""
        if not self.conversation_history:
            return "No previous conversation to follow up on. Please ask a question first."
        try:
            # Get recent conversation context
            recent_context = ""
            for i, conv in enumerate(self.conversation_history[-3:]):  # Last 3 exchanges
                recent_context += f"Previous Q: {conv['question']}\nPrevious A: {conv['answer']}\n\n"
            # Get document context
            results = self.vector_store.search(follow_up_question, k=k)
            document_context = "\n".join([
                f"[source: {chunk['metadata'].get('source', 'unknown')}, chunk: {chunk['metadata'].get('chunk_id', '?')}]\n{chunk['text']}"
                for chunk in results
            ]) if results else ""
            # Combine contexts
            full_context = f"Conversation History:\n{recent_context}\nDocument Context:\n{document_context}"
            # Generate follow-up answer
            answer = self.llm.generate_answer(follow_up_question, full_context, normalize_length)
            # Store in conversation history with chunk metadata from top result
            top_chunk = results[0] if results else None
            self.conversation_history.append({
                'question': follow_up_question,
                'answer': answer,
                'context': full_context,
                'timestamp': datetime.now().isoformat(),
                'source': top_chunk['metadata'].get('source') if top_chunk else None,
                'chunk_id': top_chunk['metadata'].get('chunk_id') if top_chunk else None,
                'chunk_text': top_chunk['text'] if top_chunk else None,
                'page': top_chunk['metadata'].get('page') if top_chunk and 'page' in top_chunk['metadata'] else None
            })
            return answer
        except Exception as e:
            self.error_count += 1
            logger.error(f"Error processing follow-up question: {str(e)}")
            return f"Error processing follow-up question: {str(e)}"

    def clear_conversation_history(self):
        """Clear the conversation history."""
        self.conversation_history = []

    def get_conversation_stats(self):
        """Get conversation statistics."""
        return {
            'total_questions': len(self.conversation_history),
            'error_count': self.error_count,
            'last_question_time': self.conversation_history[-1]['timestamp'] if self.conversation_history else None
        }

    def _is_french_question(self, text: str) -> bool:
        """Improved French detection that's more accurate and less prone to false positives."""
        text_lower = text.lower()
        
        # More specific French words and phrases that are less likely to appear in English
        french_indicators = [
            'est-ce que', 'qu\'est-ce que', 'comment ça va', 'bonjour', 'au revoir', 'merci',
            's\'il vous plaît', 'excusez-moi', 'je voudrais', 'je veux', 'je peux',
            'nous avons', 'ils ont', 'elles ont', 'c\'est', 'ce sont', 'il y a',
            'pourquoi', 'comment', 'quand', 'où', 'qui', 'quoi', 'quel', 'quelle',
            'combien', 'depuis quand', 'jusqu\'à quand', 'd\'accord', 'bien sûr',
            'peut-être', 'certainement', 'absolument', 'probablement', 'sûrement'
        ]
        
        # French accented characters (strong indicator)
        french_chars = ['é', 'è', 'ê', 'ë', 'à', 'â', 'ô', 'ù', 'û', 'ç', 'î', 'ï', 'œ', 'æ']
        
        # French-specific word patterns
        french_patterns = [
            r'\b(le|la|les|un|une|des|du|de|au|aux|dans|sur|avec|pour|par|sans|sous|chez)\b',
            r'\b(et|ou|mais|donc|car|ni|puis)\b',
            r'\b(je|tu|il|elle|nous|vous|ils|elles)\b',
            r'\b(mon|ma|mes|ton|ta|tes|son|sa|ses|notre|votre|leur)\b',
            r'\b(ce|cette|ces|celui|celle|ceux|celles)\b'
        ]
        
        import re
        
        # Check for French indicators
        french_score = 0
        
        # Strong indicators (French-specific words/phrases)
        for indicator in french_indicators:
            if indicator in text_lower:
                french_score += 3  # High weight for specific French phrases
        
        # French accented characters (very strong indicator)
        char_count = sum(1 for char in french_chars if char in text)
        french_score += char_count * 2  # High weight for accented characters
        
        # French word patterns (medium weight)
        for pattern in french_patterns:
            matches = re.findall(pattern, text_lower)
            french_score += len(matches) * 0.5  # Lower weight for common words
        
        # Require a higher threshold to avoid false positives
        # At least 2 strong indicators or 1 strong + multiple weak indicators
        return french_score >= 2

    def _translate_to_french(self, text: str) -> str:
        """Translate English text to French using Deep Translator with chunking for long texts."""
        try:
            from deep_translator import GoogleTranslator
            
            # Check if text is too long for single translation
            MAX_CHARS_PER_CHUNK = 4000  # Safe limit below 5000
            
            if len(text) <= MAX_CHARS_PER_CHUNK:
                # Short text - translate directly
                translator = GoogleTranslator(source='en', target='fr')
                result = translator.translate(text)
                return result
            else:
                # Long text - split into chunks and translate separately
                logger.info(f"Text too long ({len(text)} chars), translating in chunks...")
                
                # Split text into sentences to avoid breaking mid-sentence
                sentences = text.split('. ')
                chunks = []
                current_chunk = ""
                
                for sentence in sentences:
                    # Add period back if it's not the last sentence
                    if sentence != sentences[-1]:
                        sentence += '. '
                    
                    # Check if adding this sentence would exceed the limit
                    if len(current_chunk) + len(sentence) <= MAX_CHARS_PER_CHUNK:
                        current_chunk += sentence
                    else:
                        # Current chunk is full, save it and start a new one
                        if current_chunk:
                            chunks.append(current_chunk.strip())
                        current_chunk = sentence
                
                # Add the last chunk if it has content
                if current_chunk.strip():
                    chunks.append(current_chunk.strip())
                
                logger.info(f"Split text into {len(chunks)} chunks for translation")
                
                # Translate each chunk
                translator = GoogleTranslator(source='en', target='fr')
                translated_chunks = []
                
                for i, chunk in enumerate(chunks):
                    try:
                        logger.info(f"Translating chunk {i+1}/{len(chunks)} ({len(chunk)} chars)")
                        translated_chunk = translator.translate(chunk)
                        translated_chunks.append(translated_chunk)
                    except Exception as e:
                        logger.error(f"Error translating chunk {i+1}: {str(e)}")
                        # If translation fails for a chunk, keep the original
                        translated_chunks.append(chunk)
                
                # Combine translated chunks
                final_translation = ' '.join(translated_chunks)
                logger.info(f"Translation completed: {len(final_translation)} characters")
                return final_translation

        except Exception as e:
            logger.error(f"Error translating to French: {str(e)}")
            return text  # Return original text if translation fails

### =================== Session Manager =================== ###
class SessionManager:
    def __init__(self):
        self.sessions = {}
        self.session_timeout = 3600  # 1 hour timeout
        self.cleanup_interval = 300  # Clean up every 5 minutes
        self.last_cleanup = time.time()

    def create_session(self, session_id: str) -> Dict:
        """Create a new session."""
        self.sessions[session_id] = {
            'created_at': time.time(),
            'last_activity': time.time(),
            'rag_system': None,
            'documents_loaded': False,
            'conversation_history': []
        }
        return self.sessions[session_id]

    def get_session(self, session_id: str) -> Optional[Dict]:
        """Get an existing session."""
        self._cleanup_expired_sessions()
        
        if session_id in self.sessions:
            self.sessions[session_id]['last_activity'] = time.time()
            return self.sessions[session_id]
        return None

    def update_session(self, session_id: str, updates: Dict):
        """Update session data."""
        if session_id in self.sessions:
            self.sessions[session_id].update(updates)
            self.sessions[session_id]['last_activity'] = time.time()

    def delete_session(self, session_id: str):
        """Delete a session."""
        if session_id in self.sessions:
            del self.sessions[session_id]

    def _cleanup_expired_sessions(self):
        """Clean up expired sessions."""
        current_time = time.time()
        if current_time - self.last_cleanup > self.cleanup_interval:
            expired_sessions = []
            for session_id, session_data in self.sessions.items():
                if current_time - session_data['last_activity'] > self.session_timeout:
                    expired_sessions.append(session_id)
            
            for session_id in expired_sessions:
                del self.sessions[session_id]
            
            self.last_cleanup = current_time

### =================== Main RAG System =================== ###
class RAGSystem:
    def __init__(self, settings=None, is_web=False, use_vision_api=True, session_id: str = None, internet_mode=False):
        self.file_handler = WebFileHandler() if is_web else None
        self.vector_store = VectorStore()
        self.question_handler = QuestionHandler(self.vector_store)
        self.is_web = is_web
        self.conversation_history = []
        self.session_id = session_id
        self.session_manager = SessionManager()
        self.error_handling_enabled = True
        self.internet_mode = internet_mode
        self.image_processor = OnDemandImageProcessor()
        self.performance_metrics = {
            'total_queries': 0,
            'avg_response_time': 0,
            'error_count': 0
        }
        self.chart_extractor = PDFChartExtractor()  # Add this line

    def process_web_uploads(self, uploaded_files):
        """Process multiple uploaded files."""
        try:
            if not uploaded_files:
                return False
            
            all_documents = []
            
            for uploaded_file in uploaded_files:
                documents = self.file_handler._process_single_file(uploaded_file)
                if documents:
                    all_documents.extend(documents)
            
            if all_documents:
                # Add documents to vector store
                self.vector_store.add_documents(all_documents)
                return True
            else:
                return False
                
        except Exception as e:
            logger.error(f"Error processing web uploads: {str(e)}")
            return False

    def get_image_path(self, page_num, image_num=None):
        """Get image path for display - delegates to TextProcessor."""
        if hasattr(self.file_handler, 'text_processor'):
            return self.file_handler.text_processor.get_image_path(page_num, image_num)
        return None

    def get_all_images(self):
        """Get all extracted images info."""
        if hasattr(self.file_handler, 'text_processor'):
            return self.file_handler.text_processor.get_all_images()
        return {}

    def add_to_conversation_history(self, question, answer, question_type="initial", mode="document", source=None, page=None, chunk_text=None):
        """Add to conversation history with mode tracking."""
        entry = {
            'question': question,
            'answer': answer,
            'question_type': question_type,
            'mode': mode,  # Track which mode was used
            'timestamp': datetime.now().isoformat()
        }
        
        # Add chunk metadata if provided (for source requests)
        if source:
            entry['source'] = source
        if page:
            entry['page'] = page
        if chunk_text:
            entry['chunk_text'] = chunk_text
            
        self.conversation_history.append(entry)

    def get_conversation_history(self):
        """Get the conversation history"""
        return self.conversation_history

    def clear_conversation_history(self):
        """Clear the conversation history"""
        self.conversation_history = []

    def get_performance_metrics(self):
        """Get performance metrics for monitoring."""
        return self.performance_metrics

    def reset_performance_metrics(self):
        """Reset performance metrics."""
        self.performance_metrics = {
            'total_queries': 0,
            'avg_response_time': 0,
            'error_count': 0
        }

    def set_internet_mode(self, enabled: bool):
        """Set internet mode on or off."""
        self.internet_mode = enabled

    def is_internet_mode_enabled(self) -> bool:
        """Check if internet mode is enabled."""
        return self.internet_mode

    def generate_internet_answer(self, question: str) -> str:
        """Generate answer using internet search."""
        try:
            # Use the ClaudeHandler with internet mode system prompt
            internet_llm = ClaudeHandler(
                system_prompt="You are a financial consultant with access to live web information. Answer questions based on the provided web search results. CRITICAL: When asked about financial data, revenue, projections, or specific figures, ALWAYS lead with the most specific and detailed numerical information available. Do NOT start with vague summaries or general statements. If the context contains specific numbers, dates, or financial figures, present those FIRST. Only provide general commentary AFTER presenting the specific data. Always cite your sources and provide accurate, up-to-date information from the web. No one sentence answers. Always provide a detailed answer."
            )
            
            # For now, return a simple response indicating internet mode
            return f"Internet mode is enabled. This would search the web for: {question}"
            
        except Exception as e:
            logger.error(f"Error generating internet answer: {str(e)}")
            return f"Error generating internet answer: {str(e)}"

    def process_question_with_mode(self, question: str, normalize_length: bool = True, analysis_mode: str = "General") -> str:
        """Process question using either document mode or internet mode, supporting analysis_mode."""
        if self.internet_mode:
            logger.info("Processing question using internet mode")
            answer = generate_live_web_answer(question)
            self.add_to_conversation_history(question, answer, "internet")
            return answer
        else:
            if not self.vector_store.is_ready():
                answer = "No documents loaded. Please upload documents first or enable internet mode."
                self.add_to_conversation_history(question, answer, "error", "document")
                return answer
            logger.info(f"Processing question using document mode with analysis_mode={analysis_mode}")
            answer = self.question_handler.process_question(question, normalize_length=normalize_length, analysis_mode=analysis_mode)
            self.add_to_conversation_history(question, answer, "document")
            return answer

    def process_follow_up_with_mode(self, follow_up_question: str, normalize_length: bool = True) -> str:
        # Get conversation history
        history = self.get_conversation_history()
        # Find the last image analysis in the history
        last_image_analysis = None
        for msg in reversed(history):
            if msg.get('question_type') == 'image_analysis':
                last_image_analysis = msg['answer']
                break

        # Build context for the LLM
        if last_image_analysis:
            system_prompt = (
                "You are a helpful assistant. The user previously uploaded an image and you analyzed it. "
                "Use the following image analysis as context for the user's follow-up question. "
                "Be specific and only use the information from the analysis and the follow-up question."
            )
            user_prompt = (
                f"Previous image analysis:\n{last_image_analysis}\n\n"
                f"Follow-up question: {follow_up_question}"
            )
        else:
            system_prompt = "You are a helpful assistant. Answer the user's question based on the conversation so far."
            user_prompt = follow_up_question

        # Call OpenAI API using the newer client approach
        try:
            client = openai.OpenAI(api_key=OPENAI_API_KEY)
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ],
                max_tokens=600,
                temperature=0.3
            )
            answer = response.choices[0].message.content.strip()
        except Exception as e:
            answer = f"Error generating answer: {str(e)}"

        # Add to conversation history as a follow-up
        self.add_to_conversation_history(
            follow_up_question,
            answer,
            "image_followup" if last_image_analysis else "document_followup",
            "image" if last_image_analysis else "document"
        )
        return answer

    def get_mode_status(self) -> Dict:
        """Get current mode status and information."""
        return {
            'internet_mode': self.internet_mode,
            'documents_loaded': self.vector_store.is_ready(),
            'total_documents': len(self.vector_store.documents) if self.vector_store.documents else 0,
            'mode_description': 'Internet Search' if self.internet_mode else 'Document Search'
        }

    def handle_follow_up(self, follow_up_question: str, normalize_length: bool = True, analysis_mode: str = "General"):
        import time
        start_time = time.time()
        try:
            answer = self.process_follow_up_with_mode(follow_up_question, normalize_length=normalize_length, analysis_mode=analysis_mode)
            response_time = time.time() - start_time
            if not hasattr(self, 'performance_metrics'):
                self.performance_metrics = {}
            self.performance_metrics['last_response_time'] = response_time
            return answer
        except Exception as e:
            if hasattr(self, 'performance_metrics'):
                self.performance_metrics['error_count'] = self.performance_metrics.get('error_count', 0) + 1
            import logging
            logging.getLogger(__name__).error(f"Error generating follow-up: {str(e)}")
            return f"Error: {str(e)}"

    def process_live_web_question(self, question: str, analysis_mode: str = "General") -> str:
        """Process question using live web search."""
        try:
            # Use the analysis_mode to customize the web search prompt
            if analysis_mode == "Financial Document":
                enhanced_question = f"Financial analysis: {question}"
            elif analysis_mode == "Company Evaluation":
                enhanced_question = f"Company evaluation: {question}"
            elif analysis_mode == "Legal Document":
                enhanced_question = f"Legal analysis: {question}"
            elif analysis_mode == "Financial Excel Document":
                enhanced_question = f"Financial data analysis: {question}"
            else:
                enhanced_question = question
            
            answer = generate_live_web_answer(enhanced_question)
            self.add_to_conversation_history(question, answer, "live_web_search", "internet")
            return answer
        except Exception as e:
            logger.error(f"Error processing live web question: {str(e)}")
            return f"Error: {str(e)}"

    def process_image_request(self, question: str, pdf_path: str = None) -> str:
        """Process a request for images from a document, focusing on relevant pages."""
        try:
            logger.info(f"Processing image request: {question}")
            
            # Determine which PDF to search
            if pdf_path is None:
                # Use the most recent document from conversation history
                conversation_history = self.get_conversation_history()
                for conv in reversed(conversation_history):
                    if conv.get('source'):
                        pdf_path = conv['source']
                        break
                
                if not pdf_path:
                    return "No document specified. Please upload a document first."
            
            # Check if PDF exists
            if not os.path.exists(pdf_path):
                # Try temp directory
                temp_path = os.path.join("temp", os.path.basename(pdf_path))
                if os.path.exists(temp_path):
                    pdf_path = temp_path
                else:
                    return f"Document not found: {pdf_path}"
            
            # Get relevant pages from recent conversation or search results
            relevant_pages = self._get_relevant_pages_for_question(question, pdf_path)
            
            if not relevant_pages:
                # Fallback to processing all pages if no relevant pages found
                logger.info("No specific relevant pages found, processing entire document")
                all_images = self.image_processor.extract_images_from_pdf(pdf_path)
            else:
                # Process only relevant pages
                logger.info(f"Processing images from relevant pages: {relevant_pages}")
                all_images = self.image_processor.extract_images_from_relevant_pages(pdf_path, relevant_pages)
            
            if not all_images:
                return "No images found in the relevant pages of the document."
            
            logger.info(f"Found {len(all_images)} images, analyzing...")
            
            # ANALYZE EACH IMAGE
            results = []
            for i, img_info in enumerate(all_images):
                try:
                    logger.info(f"Analyzing image {i+1}/{len(all_images)}: {img_info['path']}")
                    analysis = self.image_processor.analyze_image_with_gpt4(img_info['path'], question)
                    
                    results.append({
                        **img_info,
                        'analysis': analysis
                    })
                    
                except Exception as e:
                    logger.error(f"Error analyzing image {img_info['path']}: {str(e)}")
                    results.append({
                        **img_info,
                        'analysis': f"Error analyzing this image: {str(e)}"
                    })
            
            # FORMAT THE RESPONSE
            response_parts = []
            response_parts.append(f"📊 **Found {len(results)} image(s) in relevant pages:**")
            response_parts.append("")
            
            for i, img_info in enumerate(results, 1):
                response_parts.append(f"**Image {i} (Page {img_info['page']}):**")
                response_parts.append(img_info['analysis'])
                response_parts.append("")
            
            return "\n".join(response_parts)
            
        except Exception as e:
            logger.error(f"Error processing image request: {str(e)}")
            return f"Error processing image request: {str(e)}"

    def process_chart_request(self, question: str, pdf_path: str = None) -> list:
        """Process requests for charts/graphs and return image files."""
        try:
            logger.info(f"Starting chart request processing for: {question}")
            
            # Determine which PDF to search
            if pdf_path is None:
                # First try to get PDF paths from the file handler
                if hasattr(self, 'file_handler') and self.file_handler:
                    saved_paths = self.file_handler.get_saved_pdf_paths()
                    if saved_paths:
                        pdf_path = saved_paths[0]
                        logger.info(f"Using PDF path from file handler: {pdf_path}")
                
                # If no path from file handler, try conversation history
                if not pdf_path:
                    conversation_history = self.get_conversation_history()
                    if conversation_history:
                        for conv in reversed(conversation_history):
                            if 'source' in conv:
                                pdf_path = conv['source']
                                logger.info(f"Using PDF path from conversation: {pdf_path}")
                                break
            
            if not pdf_path:
                logger.error("No PDF path found")
                return []
            
            if not os.path.exists(pdf_path):
                logger.error(f"PDF file does not exist: {pdf_path}")
                return []
            
            logger.info(f"Processing PDF: {pdf_path}")
            
            # Check if asking for a specific page
            import re
            page_match = re.search(r'page\s+(\d+)', question.lower())
            if page_match:
                page_number = int(page_match.group(1))
                logger.info(f"Converting specific page {page_number} to image")
                image_path = self.chart_extractor.convert_single_page_to_image(pdf_path, page_number)
                if image_path and os.path.exists(image_path):
                    logger.info(f"Successfully created image: {image_path}")
                    return [{
                        'type': 'chart_image',
                        'page': page_number,
                        'image_path': image_path,
                        'description': f"Chart from page {page_number}"
                    }]
                else:
                    logger.error(f"Failed to create image for page {page_number}")
                    return []
            
            # Get relevant pages for the question
            relevant_pages = self._get_relevant_pages_for_question(question, pdf_path)
            
            if not relevant_pages:
                # If no relevant pages found, try first few pages (but check PDF length first)
                import fitz
                doc = fitz.open(pdf_path)
                total_pages = len(doc)
                doc.close()
                
                # Only use pages that actually exist
                relevant_pages = [i for i in [1, 2, 3] if i <= total_pages]
                logger.info(f"No relevant pages found, using first {len(relevant_pages)} pages: {relevant_pages}")
            
            logger.info(f"Converting pages to images: {relevant_pages}")
            
            # Convert relevant pages to images
            chart_images = []
            for page_num in relevant_pages:
                logger.info(f"Converting page {page_num} to image")
                image_path = self.chart_extractor.convert_single_page_to_image(pdf_path, page_num)
                if image_path and os.path.exists(image_path):
                    logger.info(f"Successfully created image for page {page_num}: {image_path}")
                    chart_images.append({
                        'type': 'chart_image',
                        'page': page_num,
                        'image_path': image_path,
                        'description': f"Chart from page {page_num}"
                    })
                else:
                    logger.error(f"Failed to create image for page {page_num}")
            
            logger.info(f"Generated {len(chart_images)} chart images")
            return chart_images
            
        except Exception as e:
            logger.error(f"Error processing chart request: {str(e)}")
            return []

    def _get_relevant_pages_for_question(self, question: str, pdf_path: str) -> list:
        """Get relevant page numbers for a question."""
        try:
            relevant_pages = set()
            
            # Method 1: Get pages from recent search results
            search_results = self.vector_store.search(question, k=5)
            for chunk in search_results:
                page = chunk['metadata'].get('page')
                if page:
                    relevant_pages.add(int(page))
            
            # Method 2: Get pages from recent conversation history
            conversation_history = self.get_conversation_history()
            for conv in conversation_history[-5:]:  # Last 5 conversations
                page = conv.get('page')
                if page:
                    relevant_pages.add(int(page))
            
            # Method 3: Check if question mentions specific pages
            import re
            page_matches = re.findall(r'page\s+(\d+)', question.lower())
            for page_match in page_matches:
                relevant_pages.add(int(page_match))
            
            logger.info(f"Found relevant pages: {sorted(relevant_pages)}")
            return sorted(relevant_pages)
            
        except Exception as e:
            logger.error(f"Error getting relevant pages: {str(e)}")
            return []

    # Add this debug method to the RAGSystem class (add it after the existing methods)
    def debug_chart_processing(self, question: str):
        """Debug method to see what's happening with chart processing."""
        try:
            logger.info("=== DEBUGGING CHART PROCESSING ===")
            
            # Check if file handler exists
            if hasattr(self, 'file_handler') and self.file_handler:
                saved_paths = self.file_handler.get_saved_pdf_paths()
                logger.info(f"File handler saved paths: {saved_paths}")
            else:
                logger.info("No file handler found")
            
            # Check conversation history
            conv_history = self.get_conversation_history()
            logger.info(f"Conversation history length: {len(conv_history)}")
            
            # Try to get PDF path
            pdf_path = None
            if hasattr(self, 'file_handler') and self.file_handler:
                saved_paths = self.file_handler.get_saved_pdf_paths()
                if saved_paths:
                    pdf_path = saved_paths[0]
                    logger.info(f"Using PDF path from file handler: {pdf_path}")
            
            if not pdf_path:
                for conv in reversed(conv_history):
                    if conv.get('source'):
                        pdf_path = conv['source']
                        logger.info(f"Using PDF path from conversation: {pdf_path}")
                        break
            
            if not pdf_path:
                logger.error("No PDF path found!")
                return "No PDF path found"
            
            if not os.path.exists(pdf_path):
                logger.error(f"PDF file does not exist: {pdf_path}")
                return f"PDF file does not exist: {pdf_path}"
            
            # Check if chart_extractor exists
            if not hasattr(self, 'chart_extractor'):
                logger.error("No chart_extractor found!")
                return "No chart_extractor found"
            
            # Try to convert first page
            logger.info("Attempting to convert page 1 to image...")
            image_path = self.chart_extractor.convert_single_page_to_image(pdf_path, 1)
            
            if image_path:
                logger.info(f"Successfully created image: {image_path}")
                if os.path.exists(image_path):
                    logger.info("Image file exists on disk")
                    return f"Success! Image created: {image_path}"
                else:
                    logger.error("Image file does not exist on disk")
                    return "Image file does not exist on disk"
            else:
                logger.error("convert_single_page_to_image returned None")
                return "convert_single_page_to_image returned None"
                
        except Exception as e:
            logger.error(f"Debug error: {str(e)}")
            return f"Debug error: {str(e)}"

    # Add this method to your RAGSystem class
    def analyze_image_with_gpt4(self, image_path: str, question: str = None) -> str:
        """Analyze image using GPT-4 Vision API with enhanced prompts."""
        try:
            import base64
            
            # Read and encode image
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Determine analysis type based on question
            if "chart" in question.lower() or "graph" in question.lower() or "data" in question.lower():
                system_prompt = """You are an expert data analyst. Analyze this chart/graph and provide:
1. **Data Extraction**: Extract all numerical values, labels, and data points
2. **Trend Analysis**: Identify key trends, patterns, and relationships
3. **Key Insights**: Highlight the most important findings
4. **Context**: Provide context for what the data represents
5. **Recommendations**: If applicable, suggest actions based on the data

Present the information in a clear, structured format with specific numbers and details."""
            
            elif "text" in question.lower() or "ocr" in question.lower():
                system_prompt = """You are an expert OCR specialist. Extract and transcribe all text from this image:
1. **Text Content**: Transcribe all visible text accurately
2. **Formatting**: Preserve formatting, structure, and layout
3. **Tables**: If tables are present, format them clearly
4. **Handwriting**: If handwritten text is present, transcribe it
5. **Special Characters**: Include all symbols, numbers, and special characters

Present the extracted text in a clean, readable format."""
            
            else:
                system_prompt = """You are an expert image analyst. Analyze this image comprehensively:
1. **Visual Content**: Describe what you see in detail
2. **Objects/People**: Identify objects, people, or elements present
3. **Context**: Provide context about the image content
4. **Text**: Extract any visible text or labels
5. **Charts/Data**: If charts or data are present, analyze them
6. **Insights**: Provide relevant insights and observations

Be thorough and specific in your analysis."""
            
            # Prepare messages
            messages = [
                {
                    "role": "system",
                    "content": system_prompt
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": question or "Please analyze this image in detail."
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_data}"
                            }
                        }
                    ]
                }
            ]
            
            # Call GPT-4 Vision API
            client = openai.OpenAI(api_key=OPENAI_API_KEY)
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=messages,
                max_tokens=4096,
                temperature=0.3
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            logger.error(f"Error in image analysis: {str(e)}")
            return f"Error analyzing image: {str(e)}"

# --- Add helper functions for text extraction ---

def render_pdf_page_with_highlight(pdf_path, page_num, highlight_text=None):
    """Render a PDF page as an image, optionally highlighting the given text."""
    doc = fitz.open(pdf_path)
    page = doc.load_page(page_num - 1)  # 0-based index
    if highlight_text:
        text_instances = page.search_for(highlight_text)
        for inst in text_instances:
            page.add_highlight_annot(inst)
    pix = page.get_pixmap(dpi=200)
    img_path = f"temp/page_{page_num}_highlighted.png" if highlight_text else f"temp/page_{page_num}_screenshot.png"
    pix.save(img_path)
    doc.close()
    return img_path

def render_chunk_source_image(source_path, page_num, chunk_text):
    """Render a PDF page as an image, highlighting the chunk text if possible."""
    import fitz
    import os
    os.makedirs("temp", exist_ok=True)
    doc = fitz.open(source_path)
    page = doc.load_page(page_num - 1)  # 0-based index
    # Try to highlight all instances of the chunk text
    if chunk_text:
        text_instances = page.search_for(chunk_text)
        for inst in text_instances:
            page.add_highlight_annot(inst)
    pix = page.get_pixmap(dpi=200)
    img_path = f"temp/page_{page_num}_chunk_highlighted.png"
    pix.save(img_path)
    doc.close()
    return img_path

def batch_documents_by_token_limit(documents, max_tokens=16384):
    enc = tiktoken.get_encoding("cl100k_base")
    batches = []
    current_batch = []
    current_tokens = 0
    for doc in documents:
        tokens = len(enc.encode(doc['text']))
        if tokens > max_tokens:
            print(f"SKIPPING CHUNK: {tokens} tokens (too large for a single batch)")
            continue  # Skip this chunk
        if current_tokens + tokens > max_tokens and current_batch:
            print(f"Creating batch with {len(current_batch)} docs, {current_tokens} tokens")
            batches.append(current_batch)
            current_batch = []
            current_tokens = 0
        current_batch.append(doc)
        current_tokens += tokens
    if current_batch:
        print(f"Creating batch with {len(current_batch)} docs, {current_tokens} tokens")
        batches.append(current_batch)
    return batches

def search_web_live(query: str, num_results: int = 5) -> list:
    """Search the web using DuckDuckGo for live results."""
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=num_results))
        
        web_results = []
        for result in results:
            web_results.append({
                "title": result.get("title", ""),
                "snippet": result.get("body", ""),
                "link": result.get("link", ""),
                "source": "DuckDuckGo Live Search"
            })
        
        return web_results
    except Exception as e:
        logger.error(f"DuckDuckGo search error: {str(e)}")
        return []

def generate_live_web_answer(question: str) -> str:
    """Generate answer using live web search results."""
    try:
        # Search the web live
        web_results = search_web_live(question, num_results=5)
        
        if not web_results:
            return "No live web results found for your question."
        
        # Format web results for context
        web_context = "\n\n".join([
            f"Source: {result['title']}\n{result['snippet']}\nURL: {result['link']}"
            for result in web_results
        ])
        
        # Use OpenAI to generate answer from live web results
        client = openai.OpenAI(api_key=OPENAI_API_KEY)
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are a financial consultant with access to live web information. Answer questions based on the provided web search results. CRITICAL: When asked about financial data, revenue, projections, or specific figures, ALWAYS lead with the most specific and detailed numerical information available. Do NOT start with vague summaries or general statements. If the context contains specific numbers, dates, or financial figures, present those FIRST. Only provide general commentary AFTER presenting the specific data. Always cite your sources and provide accurate, up-to-date information from the web. No one sentence answers. Always provide a detailed answer."
                },
                {
                    "role": "user",
                    "content": f"Live Web Search Results:\n{web_context}\n\nQuestion: {question}\n\nPlease answer based on the live web results above. Include source citations."
                }
            ],
            max_tokens=16384,
            temperature=0.3
        )
        
        return response.choices[0].message.content
        
    except Exception as e:
        logger.error(f"Error generating live web answer: {str(e)}")
        return f"Error accessing live web: {str(e)}"

# Add this new class for on-demand image processing
class OnDemandImageProcessor:
    def __init__(self):
        self.processed_images = {}
        
    def extract_images_from_pdf(self, pdf_path: str) -> list:
        """Extract images from PDF using PyMuPDF."""
        try:
            if not os.path.exists(pdf_path):
                logger.error(f"PDF file not found: {pdf_path}")
                return []
            
            doc = fitz.open(pdf_path)
            images = []
            
            for page_idx in range(len(doc)):
                try:
                    page = doc.load_page(page_idx)
                    image_list = page.get_images()
                    
                    for img_index, img in enumerate(image_list):
                        try:
                            # Get image data
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            
                            # Convert to PNG format
                            if pix.n - pix.alpha < 4:  # GRAY or RGB
                                img_data = pix.tobytes("png")
                            else:  # CMYK: convert to RGB first
                                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                                img_data = pix1.tobytes("png")
                                pix1 = None
                            
                            # Save image
                            img_filename = f"page_{page_idx + 1}_image_{img_index + 1}.png"
                            img_path = os.path.join("temp", img_filename)
                            os.makedirs("temp", exist_ok=True)
                            
                            with open(img_path, "wb") as img_file:
                                img_file.write(img_data)
                            
                            images.append({
                                'path': img_path,
                                'page': page_idx + 1,
                                'image_index': img_index + 1,
                                'description': f"Image from page {page_idx + 1}"
                            })
                            
                            pix = None
                            
                        except Exception as e:
                            logger.warning(f"Error extracting image {img_index} from page {page_idx + 1}: {str(e)}")
                            continue
                            
                except Exception as e:
                    logger.warning(f"Error processing page {page_idx + 1}: {str(e)}")
                    continue
            
            doc.close()
            logger.info(f"Extracted {len(images)} images from PDF")
            return images
        
        except Exception as e:
            logger.error(f"Error extracting images from PDF: {str(e)}")
            return []
    
    def analyze_image_with_gpt4(self, image_path: str, question: str = None) -> str:
        """Analyze an image using GPT-4 Vision API."""
        try:
            if not os.path.exists(image_path):
                return "Image file not found."
            
            # Check file size (OpenAI has limits)
            file_size = os.path.getsize(image_path)
            if file_size > 20 * 1024 * 1024:  # 20MB limit
                return "Image file too large for analysis."
            
            # Read and encode the image
            with open(image_path, "rb") as img_file:
                img_data = img_file.read()
                encoded_image = base64.b64encode(img_data).decode('ascii')
            
            # Prepare the prompt
            if question:
                system_prompt = f"Analyze this image and answer the specific question: {question}. If you see charts, graphs, or tables, extract the data and present it in a clear format."
            else:
                system_prompt = "Analyze this image and extract all text and data from any charts, graphs, or tables. Present chart data in tabular format for easy understanding."
            
            # API call to GPT-4 Vision
            headers = {
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json"
            }
            
            payload = {
                "model": "gpt-4o",
                "messages": [
                    {
                        "role": "system",
                        "content": system_prompt
                    },
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{encoded_image}"
                                }
                            }
                        ]
                    }
                ],
                "max_tokens": 4096,
                "temperature": 0.3
            }
            
            response = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers=headers,
                json=payload,
                timeout=30
            )
            
            if response.status_code == 200:
                result = response.json()
                return result["choices"][0]["message"]["content"]
            else:
                logger.error(f"GPT-4 API error: {response.status_code} - {response.text}")
                return f"Error analyzing image: API returned {response.status_code}"
                
        except Exception as e:
            logger.error(f"Error analyzing image with GPT-4: {str(e)}")
            return f"Error analyzing image: {str(e)}"

    # Add this new method to the OnDemandImageProcessor class
    def extract_images_from_relevant_pages(self, pdf_path: str, relevant_pages: list, buffer_pages: int = 2) -> list:
        """Extract images only from pages near the relevant chunks."""
        try:
            if not os.path.exists(pdf_path):
                logger.error(f"PDF file not found: {pdf_path}")
                return []
            
            doc = fitz.open(pdf_path)
            total_pages = len(doc)
            images = []
            
            # Expand relevant pages with buffer
            pages_to_process = set()
            for page in relevant_pages:
                # Add the page itself and buffer pages around it
                for i in range(max(1, page - buffer_pages), min(total_pages + 1, page + buffer_pages + 1)):
                    pages_to_process.add(i)
            
            logger.info(f"Processing images from pages: {sorted(pages_to_process)} (out of {total_pages} total pages)")
            
            for page_num in sorted(pages_to_process):
                try:
                    page_idx = page_num - 1  # Convert to 0-based index
                    page = doc.load_page(page_idx)
                    image_list = page.get_images()
                    
                    for img_index, img in enumerate(image_list):
                        try:
                            # Get image data
                            xref = img[0]
                            pix = fitz.Pixmap(doc, xref)
                            
                            # Convert to PNG format
                            if pix.n - pix.alpha < 4:  # GRAY or RGB
                                img_data = pix.tobytes("png")
                            else:  # CMYK: convert to RGB first
                                pix1 = fitz.Pixmap(fitz.csRGB, pix)
                                img_data = pix1.tobytes("png")
                                pix1 = None
                            
                            # Save image
                            img_filename = f"page_{page_num}_image_{img_index + 1}.png"
                            img_path = os.path.join("temp", img_filename)
                            os.makedirs("temp", exist_ok=True)
                            
                            with open(img_path, "wb") as img_file:
                                img_file.write(img_data)
                            
                            images.append({
                                'path': img_path,
                                'page': page_num,
                                'image_index': img_index + 1,
                                'description': f"Image from page {page_num}"
                            })
                            
                            pix = None
                            
                        except Exception as e:
                            logger.warning(f"Error extracting image {img_index} from page {page_num}: {str(e)}")
                            continue
                            
                except Exception as e:
                    logger.warning(f"Error processing page {page_num}: {str(e)}")
                    continue
            
            doc.close()
            logger.info(f"Extracted {len(images)} images from {len(pages_to_process)} relevant pages")
            return images
            
        except Exception as e:
            logger.error(f"Error extracting images from relevant pages: {str(e)}")
            return []

# Add this new class for PDF to image conversion and chart extraction
class PDFChartExtractor:
    def __init__(self):
        self.output_dir = "chart_extractions"
        os.makedirs(self.output_dir, exist_ok=True)
    
    def convert_single_page_to_image(self, pdf_path, page_number):
        """Convert a single PDF page to an image."""
        try:
            import fitz  # PyMuPDF
            
            # Open the PDF
            doc = fitz.open(pdf_path)
            
            if page_number > len(doc) or page_number < 1:
                logger.error(f"Page {page_number} not found. PDF has {len(doc)} pages.")
                return None
            
            # Get the specific page
            page = doc.load_page(page_number - 1)  # 0-based index
            
            # Render page as image
            pix = page.get_pixmap(dpi=200)
            
            # Save the page image
            image_path = os.path.join(self.output_dir, f"page_{page_number}.png")
            pix.save(image_path)
            
            doc.close()
            
            logger.info(f"Successfully converted page {page_number} to {image_path}")
            return image_path
            
        except Exception as e:
            logger.error(f"Error converting page {page_number} to image: {str(e)}")
            return None

    def convert_pdf_to_images(self, pdf_path):
        """Convert PDF pages to images using PyMuPDF instead of pdf2image."""
        try:
            import fitz  # PyMuPDF
            
            # Open the PDF
            doc = fitz.open(pdf_path)
            image_paths = []
            
            # Convert each page to image
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Render page as image
                pix = page.get_pixmap(dpi=200)
                
                # Save image
                image_path = os.path.join(self.output_dir, f"page_{page_num + 1}.png")
                pix.save(image_path)
                image_paths.append(image_path)
            
            doc.close()
            return image_paths
            
        except Exception as e:
            logger.error(f"Error converting PDF to images with PyMuPDF: {str(e)}")
            return []

    def get_chart_data_by_page(self, pdf_path, page_number):
        """Extract chart data from a specific page using PyMuPDF."""
        try:
            import fitz  # PyMuPDF
            
            # Open the PDF
            doc = fitz.open(pdf_path)
            
            if page_number > len(doc) or page_number < 1:
                return f"Page {page_number} not found. PDF has {len(doc)} pages."
            
            # Get the specific page
            page = doc.load_page(page_number - 1)  # 0-based index
            
            # Render page as image
            pix = page.get_pixmap(dpi=200)
            
            # Save the page image
            image_path = os.path.join(self.output_dir, f"page_{page_number}.png")
            pix.save(image_path)
            
            doc.close()
            
            # Process the image
            extracted_text = self.process_image_and_extract_text(image_path)
            
            return extracted_text if extracted_text else f"No chart data found on page {page_number}"
            
        except Exception as e:
            logger.error(f"Error extracting chart data from page {page_number}: {str(e)}")
            return f"Error processing page {page_number}: {str(e)}"

    def process_image_and_extract_text(self, image_path, output_dir=None):
        """Process a single image and extract text/chart data using GPT-4 Vision."""
        try:
            if output_dir is None:
                output_dir = self.output_dir
            
            # Open and encode image
            with open(image_path, "rb") as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Use OpenAI Vision API to extract text and chart data
            client = openai.OpenAI(api_key=OPENAI_API_KEY)
            
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Extract all text and chart data from this image. If there are charts, graphs, or tables, convert them to a tabular format. Include all text content and organize chart data in a clear, structured way."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_data}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=4096,
                temperature=0.1
            )
            
            extracted_text = response.choices[0].message.content
            
            # Save extracted text to file
            page_num = os.path.basename(image_path).split('_')[1].split('.')[0]
            output_file = os.path.join(output_dir, f"page_{page_num}_extracted.txt")
            
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(extracted_text)
            
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error processing image {image_path}: {str(e)}")
            return f"Error processing image: {str(e)}"

    def process_pdf_for_charts(self, pdf_path):
        """Process entire PDF and extract chart data from all pages."""
        try:
            # Convert PDF to images
            image_paths = self.convert_pdf_to_images(pdf_path)
            
            if not image_paths:
                return {}
            
            extracted_data = {}
            
            # Process each page
            for i, image_path in enumerate(image_paths):
                page_num = i + 1
                logger.info(f"Processing page {page_num} for chart data")
                
                extracted_text = self.process_image_and_extract_text(image_path)
                if extracted_text and not extracted_text.startswith("Error"):
                    extracted_data[page_num] = extracted_text
            
            return extracted_data
            
        except Exception as e:
            logger.error(f"Error processing PDF for charts: {str(e)}")
            return {}

def convert_to_wav(input_path, output_path):
    command = [
        "ffmpeg",
        "-y",  # Overwrite output file if it exists
        "-i", input_path,
        "-ar", "16000",
        "-ac", "1",
        "-sample_fmt", "s16",
        output_path
    ]
    subprocess.run(command, check=True)

def handle_uploaded_audio(uploaded_file):
    # Save uploaded file to disk
    input_path = f"/tmp/{uploaded_file.name}"
    with open(input_path, "wb") as f:
        f.write(uploaded_file.read())
    # Convert to WAV
    output_path = input_path.rsplit(".", 1)[0] + ".wav"
    convert_to_wav(input_path, output_path)
    return output_path

if __name__ == "__main__": 
    rag_system = RAGSystem()
    rag_system.run()
